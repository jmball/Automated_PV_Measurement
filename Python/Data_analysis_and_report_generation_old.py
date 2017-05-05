# This script takes the measurement log files and loads them as Pandas
# DataFrames for manipulation and then plotting.

import argparse
import itertools
import os
from datetime import date, timedelta
from functools import reduce

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import scipy as sp
from matplotlib import axes, gridspec
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
from scipy import constants, signal

import reportgenlib as rgl

# Bind subplot formatting methods in reportgenlib to the matplotlib.axes.Axes
# class.
axes.Axes.set_axes_props = rgl.set_axes_props
axes.Axes.subboxplot = rgl.subboxplot
axes.Axes.subbarchart = rgl.subbarchart

# Parse folder path and log file name from command line arguments.
# Remember to include the "python" keyword before the call to the python file
# from the command line, e.g. python example.py "arg1" "arg2". Folder paths
# must use forward slashes to separate subfolders.
parser = argparse.ArgumentParser(description='Process data files')
parser.add_argument(
    'folder_path',
    metavar='folder_path',
    type=str,
    help='Absolute path to the folder containing all of the experimental data')
parser.add_argument('log_file_name',
                    metavar='log_file_name',
                    type=str,
                    help='Name of the master J-V log file')
args = parser.parse_args()

# Set folder/file paths for data and log files. Remember to always use forward
# slashes for paths.
folderpath = args.folder_path
folderpath_jv = 'J-V/'
folderpath_time = 'Time Dependence/'
folderpath_maxp = 'Max P Stabilisation/'
folderpath_intensity = 'Intensity Dependence/'
folderpath_eqe = 'EQE/'
filepath_jv = args.log_file_name
filepath_eqe = r'_EQE_LOG.txt'
log_file_jv = folderpath + folderpath_jv + filepath_jv
log_file_time = folderpath + folderpath_time + filepath_jv
log_file_maxp = folderpath + folderpath_maxp + filepath_jv
log_file_intensity = folderpath + folderpath_intensity + filepath_jv
log_file_eqe = folderpath + folderpath_eqe + filepath_jv

# Get username, date, and experiment title from folderpath for the title page
# of the powerpoint version of the report.
folderpath_split = folderpath.split('/')
username = folderpath_split[2]
date_title_split = folderpath_split[5].split(' ')
exp_date = date_title_split[0]
experiment_title = ' '.join(date_title_split[1:])

# Set physical constants
kB = constants.Boltzmann
q = constants.elementary_charge
c = constants.speed_of_light
h = constants.Planck
T = 300

# Initialise empty list for storing image file paths
images = []

# Create pdf for adding figures
pp = PdfPages(log_file_jv.replace('.txt', '_summary.pdf'))

# Import log file into a Pandas dataframe
data = pd.read_csv(log_file_jv,
                   delimiter='\t',
                   header=0,
                   names=['Label', 'Pixel', 'Condition', 'Variable', 'Value',
                          'Position', 'Jsc', 'Voc', 'PCE', 'FF', 'Area',
                          'Stabil_Level', 'Stabil_time', 'Meas_delay', 'Vmp',
                          'File_Path', 'Scan_rate', 'Scan_direction',
                          'Intensity'])

# Get scan number from file path for each row in log file and create new
# column indicating the scan number.
scan_num = []
for path in data['File_Path']:
    scan_i = path.find('scan', len(path) - 12)
    scan_num.append(path[scan_i:].strip('scan').strip('.txt'))
data['scan_num'] = pd.Series(scan_num, index=data.index)

# Define a function to do extra analysis on the measured JV curves
def extra_JV_analysis(filepath, Jsc, Vmp, Voc, Area):
    """

    Function for quantifying the properties of a solar cell J-V curve beyond
    the basic Jsc, Voc, FF, and PCE:

        - Estimates for the series and shunt resistances from the gradients at
        open-circuit and short-circuit, respectively.

    """

    # Import the data
    JV = np.genfromtxt(filepath, delimiter='\t')
    JV = JV[~np.isnan(JV).any(axis=1)]

    # Get indices of Jsc, Vmp, and Voc in data array
    V = np.around(JV[:, 0], 2)
    Voc = np.around(Voc, 2)
    Vmp = np.around(Vmp, 2)
    Jsc_i = np.where(V == 0)
    if Vmp > 0:
        Voc_i = np.where(V == Voc)
    else:
        Voc_i = np.where(V == -Voc)

    # Convert current density (in mA/cm^2) to current (in A) for diode
    # equivalent circuit fitting
    I = JV[:, 1] * Area / 1000
    IV = np.column_stack((JV[:, 0], I))

    # Estimate Rs and Rsh from gradients
    dI_dV = np.gradient(IV[:, 1], IV[1, 0] - IV[0, 0])
    dI_dV_sgfilter = signal.savgol_filter(dI_dV, 15, 3)
    try:
        Rs_grad = 1 / dI_dV_sgfilter[Voc_i[0][0]]
        Rsh_grad = 1 / dI_dV_sgfilter[Jsc_i[0][0]]
    except IndexError:
        Rs_grad = 0.1
        Rsh_grad = 0.1

# if np.size(Rs_grad) == 0:
#    Rs_grad = np.array([0])
#    Rsh_grad = np.array([0])

    return [Rs_grad, Rsh_grad]

# Use the extra analysis function to perform extra analysis and create new
# series for the dataframe
Rs_grad = []
Rsh_grad = []
for i in range(len(data['File_Path'])):
    file = data['File_Path'][i]
    new_path = file.replace('\\', '\\\\')
    Jsc = data['Jsc'][i]
    Vmp = data['Vmp'][i]
    Voc = data['Voc'][i]
    Area = data['Area'][i]
    if data['Condition'][i] == 'Light':
        extra_parameters = extra_JV_analysis(new_path, Jsc, Vmp, Voc, Area)
        Rs_grad.append(extra_parameters[0])
        Rsh_grad.append(extra_parameters[1])
    else:
        Rs_grad.append(0)
        Rsh_grad.append(0)

# Add new series to the dataframe
data['Rs_grad'] = pd.Series(Rs_grad, index=data.index)
data['Rsh_grad'] = pd.Series(Rsh_grad, index=data.index)

# Save the datafram containing additional analysis as a text file
data.to_csv(log_file_jv.replace('.txt', '_extra.txt'), sep='\t')

# Sort and filter the data to remove cells that are not working, leaving only
# best scan for each pixel
sorted_data = data.sort_values(['Variable', 'Value', 'Label', 'Pixel', 'PCE'],
                               ascending=[True, True, True, True, False])
filtered_data = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01)]
filtered_data_HL = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01) &
                               (sorted_data.Scan_direction == 'HL')]
filtered_data_LH = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01) &
                               (sorted_data.Scan_direction == 'LH')]
filtered_data = filtered_data.drop_duplicates(['Label', 'Pixel'])
filtered_data_HL = filtered_data_HL.drop_duplicates(['Label', 'Pixel'])
filtered_data_LH = filtered_data_LH.drop_duplicates(['Label', 'Pixel'])

# Drop pixels only working in one scan direction.
# First get the inner merge of the Label and Pixel columns, i.e. drop rows
# where label and pixel combination only occurs in one scan direction.
filtered_data_HL_t = filtered_data_HL[['Label', 'Pixel']].merge(
    filtered_data_LH[['Label', 'Pixel']],
    on=['Label', 'Pixel'],
    how='inner')
filtered_data_LH_t = filtered_data_LH[['Label', 'Pixel']].merge(
    filtered_data_HL[['Label', 'Pixel']],
    on=['Label', 'Pixel'],
    how='inner')

# Then perform inner merge of full filtered data frames with the merged
# label and pixel dataframes to get back all pixel data that works in both
# scan directions
filtered_data_HL = filtered_data_HL.merge(filtered_data_HL_t,
                                          on=['Label', 'Pixel'],
                                          how='inner')
filtered_data_LH = filtered_data_LH.merge(filtered_data_LH_t,
                                          on=['Label', 'Pixel'],
                                          how='inner')

# Calculate proportion of working pixels.
# Create groups of all and working pixels from complete and filtered
# dataframes.
group_var_s = sorted_data.drop_duplicates(['Label', 'Pixel'])
group_var_s = group_var_s.groupby(['Variable'])
group_var_f = filtered_data_HL.drop_duplicates(['Label', 'Pixel'])
group_var_f = group_var_f.groupby(['Variable'])

# For each variable, if there are any working pixels for that variable
# calculate the yield for each value of the variable if there are any working
# pixels for that value. Else add zeros. Also, make a list of names
# for each variable to use as the x-axis on bar charts.
yields_var = []
names_yield_var = []
for var_key in list(group_var_s.groups.keys()):
    group_val_s = group_var_s.get_group(var_key).groupby(['Value'])
    if var_key in list(group_var_f.groups.keys()):
        group_val_f = group_var_f.get_group(var_key).groupby(['Value'])
        yields_val = []
        names_yield_val = []
        for val_key in list(group_val_s.groups.keys()):
            names_yield_val.append(val_key)
            if val_key in list(group_val_f.groups.keys()):
                yields_val.append(len(group_val_f.get_group(val_key)) * 100 /
                                  8)
            else:
                yields_val.append(0)
        yields_var.append(yields_val)
        names_yield_var.append(names_yield_val)
    else:
        yields_var.append([0] * len(group_val_s))
        names_yield_var.append(list(group_val_s.groups.keys()))

# For each variable, if there are any working pixels for that variable
# calculate the yield for each label if there are any working
# pixels for that label. Else add zeros. Also, make a list of names
# for each variable to use as the x-axis on bar charts.
yields_var_lab = []
names_yield_var_lab = []
for var_key in list(group_var_s.groups.keys()):
    group_lab_s = group_var_s.get_group(var_key).groupby(['Label'])
    if var_key in list(group_var_f.groups.keys()):
        group_lab_f = group_var_f.get_group(var_key).groupby(['Label'])
        yields_lab = []
        names_yield_lab = []
        for lab_key in list(group_lab_s.groups.keys()):
            names_yield_lab.append(lab_key)
            if lab_key in list(group_lab_f.groups.keys()):
                yields_lab.append(len(group_lab_f.get_group(lab_key)) * 100 /
                                  8)
            else:
                yields_lab.append(0)
        yields_var_lab.append(yields_lab)
        names_yield_var_lab.append(names_yield_lab)
    else:
        yields_var_lab.append([0] * len(group_lab_s))
        names_yield_var_lab.append(list(group_lab_s.groups.keys()))

# To generate box plots the data needs to be grouped first by variable
grouped_by_var_HL = filtered_data_HL.groupby('Variable')
grouped_by_var_LH = filtered_data_LH.groupby('Variable')


# Then it needs to be grouped by variable value. Each of these groupings is
# appended to a list that is iterated upon later to generate the plots
def boxplotdata(grouped_by_var):
    """

    Function for sorting log file data grouped by variable into series for
    boxplots.

    """

    Jsc_var = []
    Voc_var = []
    FF_var = []
    PCE_var = []
    Rs_var = []
    Rsh_var = []
    names_var = []
    var_names = []
    for name, group in grouped_by_var:
        grouped_by_val = group.groupby('Value')
        var_names.append(name)
        Jsc = []
        Voc = []
        FF = []
        PCE = []
        Rs = []
        Rsh = []
        names_val = []
        for name, group in grouped_by_val:
            names_val.append(name)
            Jsc.append(np.array(group['Jsc']))
            Voc.append(np.array(group['Voc']))
            FF.append(np.array(group['FF']))
            PCE.append(np.array(group['PCE']))
            Rs.append(np.array(group['Rs_grad']))
            Rsh.append(np.array(group['Rsh_grad']))
        Jsc_var.append(Jsc)
        Voc_var.append(Voc)
        FF_var.append(FF)
        PCE_var.append(PCE)
        Rs_var.append(Rs)
        Rsh_var.append(Rsh)
        names_var.append(names_val)
    return {'Jsc_var': Jsc_var,
            'Voc_var': Voc_var,
            'FF_var': FF_var,
            'PCE_var': PCE_var,
            'Rs_var': Rs_var,
            'Rsh_var': Rsh_var,
            'names_var': names_var,
            'var_names': var_names}

# Create variables holding a dictionary for accessing lists of data for
# boxplots
boxplotdata_HL = boxplotdata(grouped_by_var_HL)
boxplotdata_HL['Yield_var'] = yields_var
boxplotdata_HL['Yield_var_lab'] = yields_var_lab
boxplotdata_LH = boxplotdata(grouped_by_var_LH)

# Define some properties of the box plots
boxprops_HL = dict(color='blue')
boxprops_LH = dict(color='red')
capprops_HL = dict(color='blue')
capprops_LH = dict(color='red')
whiskerprops_HL = dict(color='blue', linestyle='-')
whiskerprops_LH = dict(color='red', linestyle='-')
medianprops_HL = dict(color='blue')
medianprops_LH = dict(color='red')


def subboxplot(boxplotdata_HL, boxplotdata_LH, p, unit, i):
    """

    Make boxplot subplot for a given J-V curve parameter.

    - p = parameter as string e.g. 'Jsc', 'Voc', ...etc.
    - unit = unit for the parameter  as string e.g. '(mA/cm^2)' etc. The unit
    should include brackets to ensure the axis label makes sense.

    """

    plt.boxplot(boxplotdata_HL[p + '_var'][i], showfliers=False,
                labels=boxplotdata_HL['names_var'][i], boxprops=boxprops_HL,
                whiskerprops=whiskerprops_HL, capprops=capprops_HL,
                medianprops=medianprops_HL)
    plt.boxplot(boxplotdata_LH[p + '_var'][i], showfliers=False,
                labels=boxplotdata_LH['names_var'][i], boxprops=boxprops_LH,
                whiskerprops=whiskerprops_LH, capprops=capprops_LH,
                medianprops=medianprops_LH)
    plt.xticks(rotation=45, ha='right')
    plt.ylabel(p + ' ' + unit)
    plt.scatter(x,
                np.concatenate(boxplotdata_HL[p + '_var'][i]),
                c='blue',
                marker='x',
                label='H->L')
    plt.scatter(x,
                np.concatenate(boxplotdata_LH[p + '_var'][i]),
                c='red',
                marker='x',
                label='L->H')
    plt.legend(loc='lower right', scatterpoints=1, fontsize=7)
    if p == 'FF':
        plt.ylim([0, 1])
    elif (p == 'Rs') or (p == 'Rsh'):
        plt.yscale('log')
    else:
        plt.ylim(ymin=0)

# Scale figures to a landscape A4 page in inches
A4_height = 7.5
A4_width = 10

# Iterate through the lists of grouped data to produce boxplots. Each plot
# will contain all data from all values of a variable including a 'Control'
# sample if one is given. Multiple plots are created if there is more than
# one variable.
for i in range(len(boxplotdata_HL['var_names'])):
    if boxplotdata_HL['var_names'][i] != 'Control':
        for j in range(len(boxplotdata_HL['names_var'])):
            if boxplotdata_HL['names_var'][j][0] == 'Control':
                boxplotdata_HL['names_var'][i].append(boxplotdata_HL[
                    'names_var'][j][0])
                boxplotdata_HL['Jsc_var'][i].append(boxplotdata_HL['Jsc_var'][
                    j][0])
                boxplotdata_HL['Voc_var'][i].append(boxplotdata_HL['Voc_var'][
                    j][0])
                boxplotdata_HL['FF_var'][i].append(boxplotdata_HL['FF_var'][j][
                    0])
                boxplotdata_HL['PCE_var'][i].append(boxplotdata_HL['PCE_var'][
                    j][0])
                boxplotdata_HL['Rs_var'][i].append(boxplotdata_HL['Rs_var'][j][
                    0])
                boxplotdata_HL['Rsh_var'][i].append(boxplotdata_HL['Rsh_var'][
                    j][0])
        x = []
        for k in range(1, 1 + len(boxplotdata_HL['names_var'][i])):
            x.extend([k] * len(boxplotdata_HL['Jsc_var'][i][k - 1]))
        plt.figure(i + 2 * i, figsize=(A4_width, A4_height), dpi=300)
        plt.suptitle(boxplotdata_HL['var_names'][i] + ' basic parameters',
                     fontsize=14, fontweight='bold')
        plt.subplot(2, 2, 1)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'Jsc', '(mA/cm^2)', i)
        plt.subplot(2, 2, 2)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'Voc', '(V)', i)
        plt.subplot(2, 2, 3)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'FF', '', i)
        plt.subplot(2, 2, 4)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'PCE', '(%)', i)
        plt.tight_layout()
        plt.subplots_adjust(top=0.92)
        image_path = log_file_jv.replace(
            '.txt',
            '_' + boxplotdata_HL['var_names'][i] + '_basic_boxplots.png')
        plt.savefig(image_path)
        images.append(image_path)
        pp.savefig()

        plt.figure(i + 1 + 2 * i, figsize=(A4_width, A4_height), dpi=300)
        plt.suptitle(boxplotdata_HL['var_names'][i] +
                     ' series and shunt resistances; yields', fontsize=14,
                     fontweight='bold')
        plt.subplot(2, 2, 1)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'Rs', '(ohms)', i)
        plt.subplot(2, 2, 2)
        subboxplot(boxplotdata_HL, boxplotdata_LH, 'Rsh', '(ohms)', i)
        plt.subplot(2, 2, 3)
        plt.bar(
            range(len(yields_var[i])),
            yields_var[i],
            align='center',
            width=0.25,
            edgecolor='black',
            color='green')
        plt.xticks(range(len(yields_var[i])), names_yield_var[i], rotation=45)
        plt.ylabel('Yield (%)')
        plt.ylim([0, 102])
        plt.subplot(2, 2, 4)
        plt.bar(
            range(len(yields_var_lab[i])),
            yields_var_lab[i],
            align='center',
            width=0.25,
            edgecolor='black',
            color='green')
        plt.xticks(
            range(len(yields_var_lab[i])),
            names_yield_var_lab[i],
            rotation=45)
        plt.ylabel('Yield (%)')
        plt.ylim([0, 102])
        plt.tight_layout()
        plt.subplots_adjust(top=0.92)
        image_path = log_file_jv.replace(
            '.txt', '_' + boxplotdata_HL['var_names'][i] +
            '_series_shunt_resistances_yields.png')
        plt.savefig(image_path)
        images.append(image_path)
        pp.savefig()

# Group data by label and sort ready to plot graph of all pixels per substrate
re_sort_data = filtered_data.sort_values(['Label', 'Pixel'],
                                         ascending=[True, True])
grouped_by_label = re_sort_data.groupby('Label')

# Get parameters for defining position of figures in subplot, attempting to
# make it as square as possible
no_of_subplots = len(grouped_by_label)
subplot_rows = np.ceil(no_of_subplots**(1 / 2))
subplot_cols = np.ceil(no_of_subplots / subplot_rows)

# Get colormap
cmap = plt.cm.get_cmap('rainbow')

# Create lists of varibales, values, and labels for labelling figures
substrates = re_sort_data.drop_duplicates(['Label'])
variables = list(substrates['Variable'])
values = list(substrates['Value'])
labels = list(substrates['Label'])

# Create main figure
plt.figure(figsize=(A4_width, A4_height), dpi=300)
plt.suptitle('JV scans of every working pixel', fontsize=10, fontweight='bold')
i = 1
for name, group in grouped_by_label:
    plt.subplot(subplot_rows, subplot_cols, i)
    plt.axhline(0, lw=0.5, c='black')
    plt.title(
        str(labels[i - 1]) + ', ' + str(variables[i - 1]) + ', ' +
        str(values[i - 1]),
        fontsize=8)

    c_div = 1 / len(group)
    j = 0
    pixels = list(group['Pixel'])
    max_group_jsc = max(list(group['Jsc']))
    for file in group['File_Path']:
        new_path = file.replace('\\', '\\\\')

        if '_LH_' in new_path:
            data_LH_path = new_path
            data_HL_path = new_path.replace('_LH_', '_HL_')
        else:
            data_HL_path = new_path
            data_LH_path = new_path.replace('_HL_', '_LH_')

        data_LH = np.genfromtxt(data_LH_path, delimiter='\t')
        data_HL = np.genfromtxt(data_HL_path, delimiter='\t')
        data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
        data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]

        plt.plot(data_LH[:, 0],
                 data_LH[:, 1],
                 label=pixels[j],
                 c=cmap(j * c_div),
                 lw=2.0)
        plt.plot(data_HL[:, 0], data_HL[:, 1], c=cmap(j * c_div), lw=2.0)

        j += 1

    plt.xlabel('Applied bias (V)', fontsize=7)
    plt.xticks(fontsize=7)
    plt.ylabel('J (mA/cm^2)', fontsize=7)
    plt.yticks(fontsize=7)
    plt.ylim([-max_group_jsc * 1.1, max_group_jsc * 1.1])
    plt.legend(loc='best', prop={'size': 6})

    i += 1

plt.tight_layout()
plt.subplots_adjust(top=0.92)
image_path = log_file_jv.replace('.txt', '_all_JVs.png')
plt.savefig(image_path)
images.append(image_path)
pp.savefig()

# filter dataframe to leave on the best pixel for each variable value
sort_best_pixels = filtered_data.sort_values(['Variable', 'Value', 'PCE'],
                                             ascending=[True, True, False])
best_pixels = sort_best_pixels.drop_duplicates(['Variable', 'Value'])

# get parameters for defining position of figures in subplot, attempting to
# make it as square as possible
no_of_subplots = len(best_pixels['File_Path'])
subplot_rows = np.ceil(no_of_subplots**(1 / 2))
subplot_cols = np.ceil(no_of_subplots / subplot_rows)

# create lists of varibales and values for labelling figures
variables = list(best_pixels['Variable'])
values = list(best_pixels['Value'])
jscs = list(best_pixels['Jsc'])

# create main figure
plt.figure(figsize=(A4_width, A4_height), dpi=300)
plt.suptitle('Best pixels', fontsize=10, fontweight='bold')

# initialise index for subplot and start looping through best cells dataframe
i = 1

# Loop for iterating through best pixels dataframe and picking out JV data
# files. Each plot contains forward and reverse sweeps, both light and dark.
for file in best_pixels['File_Path']:
    new_path = file.replace('\\', '\\\\')

    if '_LH_' in new_path:
        JV_light_LH_path = new_path
        JV_light_HL_path = new_path.replace('_LH_', '_HL_')
    else:
        JV_light_HL_path = new_path
        JV_light_LH_path = new_path.replace('_HL_', '_LH_')

    try:
        JV_light_LH_data = np.genfromtxt(JV_light_LH_path, delimiter='\t')
        JV_light_HL_data = np.genfromtxt(JV_light_HL_path, delimiter='\t')
        JV_dark_LH_data = np.genfromtxt(
            JV_light_LH_path.replace('Light', 'Dark'),
            delimiter='\t')
        JV_dark_HL_data = np.genfromtxt(
            JV_light_HL_path.replace('Light', 'Dark'),
            delimiter='\t')
    except OSError:
        pass

    JV_light_LH_data = JV_light_LH_data[~np.isnan(JV_light_LH_data).any(
        axis=1)]
    JV_light_HL_data = JV_light_HL_data[~np.isnan(JV_light_HL_data).any(
        axis=1)]

    try:
        JV_dark_LH_data = JV_dark_LH_data[~np.isnan(JV_dark_LH_data).any(
            axis=1)]
        JV_dark_HL_data = JV_dark_HL_data[~np.isnan(JV_dark_HL_data).any(
            axis=1)]
    except NameError:
        pass

    plt.subplot(subplot_rows, subplot_cols, i)
    plt.axhline(0, lw=0.5, c='black')
    plt.title(str(variables[i - 1]) + ', ' + str(values[i - 1]), fontsize=8)
    plt.plot(JV_light_LH_data[:, 0],
             JV_light_LH_data[:, 1],
             label='L->H',
             c='red',
             lw=2.0)
    plt.plot(JV_light_HL_data[:, 0],
             JV_light_HL_data[:, 1],
             label='H->L',
             c='green',
             lw=2.0)
    try:
        plt.plot(JV_dark_LH_data[:, 0],
                 JV_dark_LH_data[:, 1],
                 label='L->H',
                 c='blue',
                 lw=2.0)
        plt.plot(JV_dark_HL_data[:, 0],
                 JV_dark_HL_data[:, 1],
                 label='H->L',
                 c='orange',
                 lw=2.0)
    except NameError:
        pass
    plt.xlabel('Applied bias (V)', fontsize=7)
    plt.xticks(fontsize=7)
    plt.ylabel('J (mA/cm^2)', fontsize=7)
    plt.yticks(fontsize=7)
    plt.ylim([-jscs[i - 1] * 1.1, jscs[i - 1] * 1.1])
    plt.legend(loc='best', prop={'size': 6})

    i += 1

plt.tight_layout()
plt.subplots_adjust(top=0.92)
image_path = log_file_jv.replace('.txt', '_best_JVs.png')
plt.savefig(image_path)
images.append(image_path)
pp.savefig()

# Sort and filter data ready for plotting different scan rates/repeat scans
sorted_data_scan = data.sort_values(
    ['Variable', 'Value', 'Label', 'Pixel', 'scan_num'],
    ascending=[True, True, True, True, True])
filtered_scan_HL = sorted_data_scan[(sorted_data_scan.Condition == 'Light') & (
    sorted_data_scan.FF > 0.1) & (sorted_data_scan.FF < 0.9) & (
        sorted_data_scan.Jsc > 0.01) & (sorted_data_scan.Scan_direction == 'HL'
                                        )]
filtered_scan_LH = sorted_data_scan[(sorted_data_scan.Condition == 'Light') & (
    sorted_data_scan.FF > 0.1) & (sorted_data_scan.FF < 0.9) & (
        sorted_data_scan.Jsc > 0.01) & (sorted_data_scan.Scan_direction == 'LH'
                                        )]

# Drop pixels only working in one scan direction
filtered_data_HL = filtered_data_HL[filtered_data_HL.Label.isin(
    filtered_data_LH.Label.values) & filtered_data_HL.Pixel.isin(
        filtered_data_LH.Pixel.values)]
filtered_data_LH = filtered_data_LH[filtered_data_LH.Label.isin(
    filtered_data_HL.Label.values) & filtered_data_LH.Pixel.isin(
        filtered_data_HL.Pixel.values)]

# Create groups of data for each pixel for a given label
group_by_label_pixel_HL = filtered_scan_HL.groupby(['Label', 'Pixel'])
group_by_label_pixel_LH = filtered_scan_LH.groupby(['Label', 'Pixel'])

# Iterate through these groups and plot JV curves if more than one scan hasn
# has been performed
i = 0
i_max = len(group_by_label_pixel_HL) - 1
j = 0
for ng_HL, ng_LH in zip(group_by_label_pixel_HL, group_by_label_pixel_LH):
    name_HL = ng_HL[0]
    group_HL = ng_HL[1]
    name_LH = ng_LH[0]
    group_LH = ng_LH[1]
    if any(int(scan) > 0 for scan in group_HL['scan_num']):
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]
        jsc_max = max(max(group_HL['Jsc']), max(group_LH['Jsc']))
        c_div = 1 / len(group_HL)
        if i % 4 == 0:
            plt.figure(figsize=(A4_width, A4_height), dpi=300)
            plt.suptitle('Repeat scan/scan rate variation JV curves',
                         fontsize=10,
                         fontweight='bold')
            j += 1
        plt.subplot(2, 2, (i % 4) + 1)
        k = 0
        for path_HL, path_LH, scan_rate_HL, scan_rate_LH, scan_num_HL, scan_num_LH in zip(
                group_HL['File_Path'], group_LH['File_Path'],
                group_HL['Scan_rate'], group_LH['Scan_rate'],
                group_HL['scan_num'], group_LH['scan_num']):
            data_HL = np.genfromtxt(path_HL, delimiter='\t')
            data_LH = np.genfromtxt(path_LH, delimiter='\t')
            data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
            data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
            plt.plot(
                data_HL[:, 0],
                data_HL[:, 1],
                c=cmap(k * c_div),
                label=str(scan_num_HL) + ', ' + str(scan_rate_HL) + ' V/s')
            plt.plot(data_LH[:, 0], data_LH[:, 1], c=cmap(k * c_div))
            k += 1
        plt.legend(loc='best', fontsize=7)
        plt.xlabel('Applied voltage (V)', fontsize=9)
        plt.ylabel('J (mA/cm^2)', fontsize=9)
        plt.ylim([-jsc_max * 1.05, jsc_max * 1.05])
        plt.title(
            str(label) + ', pixel ' + str(pixel) + ', ' + str(variable) + ', '
            + str(value),
            fontsize=8)
        if (i % 4 == 3) or (i == i_max):
            plt.tight_layout()
            plt.subplots_adjust(top=0.92)
            image_path = log_file_jv.replace(
                '.txt', '_repeat_scan_JVs_' + str(j) + '.png')
            plt.savefig(image_path)
            images.append(image_path)
            pp.savefig()
        i += 1


# Check whether other experiments have been performed and therefore if more
# data needs to be analysed.
def exp_file_list(folderpath):
    """

    Determine whether data exists for an additional experiment. If so, make
    a list of files. Returns a dictionary containing a boolean indicating
    existance and a list of files (empty if experiment hasn't been performed).

    """

    files = []
    files_split = []
    if os.path.exists(folderpath):
        exists = True
        for f in os.listdir(folderpath):
            if (f.endswith('.txt')) & (f.find('_LOG.txt') == -1):
                files.append(folderpath + f)
                files_split.append(f.split('_'))
        if len(files) == 0:
            exists = False
    else:
        exists = False

    return {'exists': exists, 'files': files, 'files_split': files_split}


# Build a log dataframe from from experiment files and master J-V log.
def build_log_df(exp_files, master_log):
    """

    Build a log dataframe from from experiment files and master J-V log.

    """

    labels = []
    pixels = []
    scan_dir = []
    for item in exp_files['files_split']:
        labels.append(item[0])
        pixels.append(item[1])
        scan_dir.append(item[3])
    files_dict = {'Label': labels,
                  'Pixel': pixels,
                  'File_Path': exp_files['files'],
                  'Scan_direction': scan_dir}
    files_df = pd.DataFrame(files_dict)
    files_df = files_df.sort_values(['Label'], ascending=True)
    unique_labels_df = master_log.drop_duplicates(['Label'])
    unique_labels_df = unique_labels_df.sort_values(['Label'], ascending=True)
    unique_labels_df = unique_labels_df[unique_labels_df.Label.isin(
        files_df.Label.values)]
    files_df['Variable'] = unique_labels_df['Variable']
    files_df['Value'] = unique_labels_df['Value']

    return files_df


# Function for calculating all factors of an integer.
def factors(n):
    """

    Returns a set of factors of the integer n.

    """

    pairs = [[i, n // i] for i in range(1, int(n**0.5) + 1) if n % i == 0]
    return set(reduce(list.__add__, pairs))


# Function for rounding a number to given number of significant figures.
def round_sig_fig(x, sf):
    """

    Rounds any number, x, to the specified number of significant figures, sf.

    """
    format_str = '%.' + str(sf) + 'e'
    x_dig, x_ord = map(float, (format_str % x).split('e'))
    return round(x, int(-x_ord) + 1)


# Function for calculting the tick positions for a plot axis.
def calc_ticks(data_min, data_max):
    """

    Returns an array of axis tick labels given the max and min values of
    the data to be plotted.

    """

    min_digs, min_order = map(float, ('%.6e' % data_min).split('e'))
    max_digs, max_order = map(float, ('%.6e' % data_max).split('e'))
    min_order -= 1
    max_order -= 1
    if max_order > min_order:
        min_digs = np.floor(min_digs * 10**(min_order - max_order + 1))
        max_digs = np.ceil(max_digs * 10)
        y_range_digs = max_digs - min_digs
        while np.max(list(filter(lambda x: x < 8, factors(y_range_digs)))) < 3:
            y_range_digs += 2
            min_digs -= 1
            max_digs += 1
        num_ticks = np.max(list(filter(lambda x: x < 8, factors(
            y_range_digs))))
        min_tick = min_digs * 10**max_order
        ticks = np.linspace(min_tick, y_range_digs * 10**max_order + min_tick,
                            num_ticks + 1)
    elif min_order > max_order:
        min_digs = np.floor(min_digs * 10)
        max_digs = np.ceil(max_digs * 10**(max_order - min_order + 1))
        y_range_digs = max_digs - min_digs
        while np.max(list(filter(lambda x: x < 8, factors(y_range_digs)))) < 3:
            y_range_digs += 2
            min_digs -= 1
            max_digs += 1
        num_ticks = np.max(list(filter(lambda x: x < 8, factors(
            y_range_digs))))
        min_tick = min_digs * 10**min_order
        ticks = np.linspace(min_tick, y_range_digs * 10**min_order + min_tick,
                            num_ticks + 1)
    elif max_order == min_order:
        str_min = str(min_digs).replace('.', '').replace('-', '')
        str_max = str(max_digs).replace('.', '').replace('-', '')
        i = 0
        while str_min[i] == str_max[i]:
            i += 1
            if i == len(str_min):
                break
        if i == len(str_min):
            ticks = np.array([min_digs * 10**(min_order + 1)])
        else:
            min_digs = np.floor(min_digs * 10**(1 + i))
            max_digs = np.ceil(max_digs * 10**(1 + i))
            y_range_digs = max_digs - min_digs
            while np.max(list(filter(lambda x: x < 8, factors(
                    y_range_digs)))) < 3:
                y_range_digs += 2
                min_digs -= 1
                max_digs += 1
            num_ticks = np.max(list(filter(lambda x: x < 8, factors(
                y_range_digs))))
            min_tick = min_digs * 10**(min_order - i)
            ticks = np.linspace(min_tick, y_range_digs * 10**
                                (min_order - i) + min_tick, num_ticks + 1)

    return ticks

# Build a time dependence log dataframe from file paths and J-V log file.
time_files = exp_file_list(folderpath + folderpath_time)
if time_files['exists']:
    time_files_df = build_log_df(time_files, sorted_data)

    sorted_time_files = time_files_df.sort_values(
        ['Label', 'Pixel', 'Scan_direction'],
        ascending=[True, True, True])
    sorted_time_files = sorted_time_files.drop_duplicates(['Label', 'Pixel'])

    i = 0
    i_max = len(sorted_time_files) - 1
    j = 0
    for index, row in sorted_time_files.iterrows():
        if i % 4 == 0:
            fig = plt.figure(figsize=(A4_width, A4_height), dpi=300)
            fig.suptitle('J-t characterics', fontsize=10, fontweight='bold')
            gs = gridspec.GridSpec(7, 2)
            j += 1
        path_HL = row['File_Path']
        path_LH = path_HL.replace('HL', 'LH')
        data_HL = np.genfromtxt(path_HL)
        data_LH = np.genfromtxt(path_LH)
        data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
        data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
        data = np.vstack([data_HL, data_LH])
        ax1 = fig.add_subplot(gs[(int(np.floor((i / 2) % 2)) * 4):(int(
            np.floor((i / 2) % 2)) * 4) + 2, i % 2])
        ax1.plot(data[:, 0], data[:, 2], 'blue')
        yticks = calc_ticks(np.min(data[:, 2]), np.max(data[:, 2]))
        ax1.set_yticks(yticks)
        ax1.set_yticklabels(yticks, fontsize=7)
        ax1.set_ylim([yticks[0], yticks[-1]])
        ax1.set_ylabel('J (mA/cm^2)', fontsize=7)
        xticks = calc_ticks(np.min(data[:, 0]), np.max(data[:, 0]))
        ax1.set_xticks(xticks)
        ax1.set_xticklabels([])
        ax1.set_xlim([xticks[0], xticks[-1]])
        ax1.set_title(
            str(row['Label']) + ', pixel ' + str(row['Pixel']) + ', ' +
            str(row['Variable']) + ', ' + str(row['Value']),
            fontsize=9)
        ax2 = fig.add_subplot(gs[(int(np.floor((i / 2) % 2)) * 4) + 2, i % 2])
        ax2.plot(data[:, 0], data[:, 1], 'black')
        yticks = calc_ticks(np.min(data[:, 1]), np.max(data[:, 1]))
        ytick_range = yticks[-1] - yticks[0]
        ax2.set_yticks(yticks)
        ax2.set_yticklabels(yticks, fontsize=7)
        ax2.set_ylim(
            [yticks[0] - 0.05 * ytick_range, yticks[-1] + 0.05 * ytick_range])
        ax2.set_ylabel('Bias (V)', fontsize=7)
        ax2.set_xticks(xticks)
        ax2.set_xticklabels(xticks, fontsize=7)
        ax2.set_xlim([xticks[0], xticks[-1]])
        ax2.set_xlabel('Time (s)', fontsize=7)
        if (i % 4 == 3) or (i == i_max):
            gs.update(wspace=0.3, hspace=0.1)
            image_path = log_file_jv.replace(
                '.txt', '_jt_characteristics' + '_' + str(j) + '.png')
            fig.savefig(image_path)
            images.append(image_path)
            pp.savefig()
        i += 1

# Build a max power stabilisation log dataframe from file paths and J-V log
# file.
maxp_files = exp_file_list(folderpath + folderpath_maxp)
if maxp_files['exists']:
    maxp_files_df = build_log_df(maxp_files, sorted_data)

    colors = ('Green', 'Red', 'Blue')
    labels = ('Applied Bias (V)', '|J| (mA/cm^2)', 'PCE (%)')
    cols = (1, 3, 5)
    i = 0
    i_max = len(maxp_files_df) - 1
    j = 0
    for index, row in maxp_files_df.iterrows():
        if i % 4 == 0:
            fig = plt.figure(figsize=(A4_width, A4_height), dpi=300)
            fig.suptitle('Maximum power stabilisation',
                         fontsize=10,
                         fontweight='bold')
            gs = gridspec.GridSpec(2, 2)
            j += 1
        path = row['File_Path']
        data = np.genfromtxt(path, delimiter='\t')
        data = data[~np.isnan(data).any(axis=1)]
        ax = fig.add_subplot(gs[int(np.floor(i / 2)), i % 2])
        ax.set_title(
            str(row['Label']) + ', pixel ' + str(row['Pixel']) + ', ' +
            str(row['Variable']) + ', ' + str(row['Value']),
            fontsize=9)
        axes = [ax, ax.twinx(), ax.twinx()]
        axes[-1].spines['right'].set_position(('axes', 1.2))
        axes[-1].set_frame_on(True)
        axes[-1].patch.set_visible(False)
        for ax, color, col, label in zip(axes, colors, cols, labels):
            ax.plot(data[:, 0], np.absolute(data[:, col]), color=color)
            ax.set_ylabel(label, color=color)
            ax.tick_params(axis='y', colors=color)
        axes[0].set_xlabel('Time (s)')
        if (i % 4 == 3) or (i == i_max):
            gs.update(wspace=0.6, hspace=0.3)
            image_path = log_file_jv.replace(
                '.txt', '_max_P_stab' + '_' + str(j) + '.png')
            fig.savefig(image_path)
            images.append(image_path)
            pp.savefig()
        i += 1

# Plot inensity dependent graphs if experiment data exists
if os.path.exists(folderpath + folderpath_intensity + filepath_jv):
    data = pd.read_csv(folderpath + folderpath_intensity + filepath_jv,
                       delimiter='\t',
                       header=0,
                       names=['Label', 'Pixel', 'Condition', 'Variable',
                              'Value', 'Position', 'Jsc', 'Voc', 'PCE', 'FF',
                              'Area', 'Stabil_Level', 'Stabil_time',
                              'Meas_delay', 'Vmp', 'File_Path', 'Scan_rate',
                              'Scan_direction', 'Intensity'])
    sorted_data_int = data.sort_values(['Label', 'Pixel', 'Intensity', 'PCE'],
                                       ascending=[True, True, True, False])
    filtered_data_int_HL = sorted_data_int[(
        sorted_data_int.Scan_direction == 'HL')]
    filtered_data_int_LH = sorted_data_int[(
        sorted_data_int.Scan_direction == 'LH')]
    filtered_data_int_HL = filtered_data_int_HL.drop_duplicates(
        ['Label', 'Pixel', 'Intensity'])
    filtered_data_int_LH = filtered_data_int_LH.drop_duplicates(
        ['Label', 'Pixel', 'Intensity'])
    group_by_label_pixel_HL = filtered_data_int_HL.groupby(['Label', 'Pixel'])
    group_by_label_pixel_LH = filtered_data_int_LH.groupby(['Label', 'Pixel'])

    # Plot intensity dependent JV curve parameters
    for ng_HL, ng_LH in zip(group_by_label_pixel_HL, group_by_label_pixel_LH):
        name_HL = ng_HL[0]
        group_HL = ng_HL[1]
        name_LH = ng_LH[0]
        group_LH = ng_LH[1]
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]
        m_HL, c_HL, r_HL, p_HL, se_HL = sp.stats.linregress(
            group_HL['Intensity'] * 100, group_HL['Jsc'])
        m_LH, c_LH, r_LH, p_LH, se_LH = sp.stats.linregress(
            group_LH['Intensity'] * 100, group_LH['Jsc'])
        r_sq_HL = r_HL**2
        r_sq_LH = r_LH**2
        plt.figure(figsize=(A4_width, A4_height), dpi=300)
        plt.suptitle('Intensity dependence, ' + str(label) + ', pixel ' +
                     str(pixel) + ', ' + str(variable) + ', ' + str(value),
                     fontsize=10,
                     fontweight='bold')
        plt.subplot(2, 2, 1)
        plt.scatter(group_HL['Intensity'] * 100,
                    group_HL['Jsc'],
                    c='blue',
                    label='H->L')
        plt.scatter(group_LH['Intensity'] * 100,
                    group_LH['Jsc'],
                    c='red',
                    label='L->H')
        plt.legend(loc='upper left', scatterpoints=1, fontsize=9)
        plt.plot(group_HL['Intensity'] * 100,
                 group_HL['Intensity'] * 100 * m_HL + c_HL,
                 c='blue')
        plt.plot(group_LH['Intensity'] * 100,
                 group_LH['Intensity'] * 100 * m_LH + c_LH,
                 c='red')
        plt.text(
            np.max(group_HL['Intensity'] * 100) * 0.6,
            (np.max(group_HL['Jsc']) - np.min(group_HL['Jsc'])) * 0.11 +
            np.min(group_HL['Jsc']),
            'm=' + str(round_sig_fig(m_HL, 3)) + ', c=' +
            str(round_sig_fig(c_HL, 3)) + ', R^2=' +
            str(round_sig_fig(r_sq_HL, 3)),
            fontsize=8,
            color='blue')
        plt.text(
            np.max(group_HL['Intensity'] * 100) * 0.6,
            (np.max(group_HL['Jsc']) - np.min(group_HL['Jsc'])) * 0.01 +
            np.min(group_HL['Jsc']),
            'm=' + str(round_sig_fig(m_LH, 3)) + ', c=' +
            str(round_sig_fig(c_LH, 3)) + ', R^2=' +
            str(round_sig_fig(r_sq_LH, 3)),
            fontsize=8,
            color='red')
        plt.xlabel('Light intensity (W/cm^2)', fontsize=9)
        plt.ylabel('Jsc (mA/cm^2)', fontsize=9)
        m_HL, c_HL, r_HL, p_HL, se_HL = sp.stats.linregress(
            np.log(group_HL['Jsc']), group_HL['Voc'])
        m_LH, c_LH, r_LH, p_LH, se_LH = sp.stats.linregress(
            np.log(group_LH['Jsc']), group_LH['Voc'])
        r_sq_HL = r_HL**2
        r_sq_LH = r_LH**2
        n_HL = m_HL * q / (kB * T)
        n_LH = m_LH * q / (kB * T)
        j0_HL = np.exp(-c_HL / m_HL)
        j0_LH = np.exp(-c_LH / m_LH)
        plt.subplot(2, 2, 2)
        plt.scatter(
            np.log(group_HL['Jsc']),
            group_HL['Voc'],
            c='blue',
            label='H->L')
        plt.scatter(
            np.log(group_LH['Jsc']),
            group_LH['Voc'],
            c='red',
            label='L->H')
        plt.legend(loc='upper left', scatterpoints=1, fontsize=9)
        plt.plot(
            np.log(group_HL['Jsc']),
            np.log(group_HL['Jsc']) * m_HL + c_HL,
            c='blue')
        plt.plot(
            np.log(group_LH['Jsc']),
            np.log(group_LH['Jsc']) * m_LH + c_LH,
            c='red')
        plt.text(
            np.max(np.log(group_HL['Jsc'])) * 0.3,
            (np.max(group_HL['Voc']) - np.min(group_HL['Voc'])) * 0.11 +
            np.min(group_HL['Voc']),
            'n=' + str(round_sig_fig(n_HL, 3)) + ', J_0=' + ('%.2e' % j0_HL) +
            ' (mA/cm^2)' + ', R^2=' + str(round_sig_fig(r_sq_HL, 3)),
            fontsize=8,
            color='blue')
        plt.text(
            np.max(np.log(group_HL['Jsc'])) * 0.3,
            (np.max(group_HL['Voc']) - np.min(group_HL['Voc'])) * 0.01 +
            np.min(group_HL['Voc']),
            'n=' + str(round_sig_fig(n_LH, 3)) + ', J_0=' + ('%.2e' % j0_LH) +
            ' (mA/cm^2)' + ', R^2=' + str(round_sig_fig(r_sq_LH, 3)),
            fontsize=8,
            color='red')
        plt.xlabel('ln(Jsc) (mA/cm^2)', fontsize=9)
        plt.ylabel('Voc (V)', fontsize=9)
        plt.subplot(2, 2, 3)
        plt.scatter(group_HL['Intensity'] * 100,
                    group_HL['FF'],
                    c='blue',
                    label='H->L')
        plt.scatter(group_LH['Intensity'] * 100,
                    group_LH['FF'],
                    c='red',
                    label='L->H')
        plt.legend(loc='lower right', scatterpoints=1, fontsize=9)
        plt.xlabel('Light intensity (W/cm^2)', fontsize=9)
        plt.ylabel('FF', fontsize=9)
        plt.subplot(2, 2, 4)
        plt.scatter(group_HL['Intensity'] * 100,
                    group_HL['PCE'],
                    c='blue',
                    label='H->L')
        plt.scatter(group_LH['Intensity'] * 100,
                    group_LH['PCE'],
                    c='red',
                    label='L->H')
        plt.legend(loc='lower right', scatterpoints=1, fontsize=9)
        plt.xlabel('Light intensity (W/cm^2)', fontsize=9)
        plt.ylabel('PCE (%)', fontsize=9)
        plt.tight_layout()
        plt.subplots_adjust(top=0.92)
        image_path = log_file_jv.replace(
            '.txt', '_intensity_' + str(label) + '_' + str(pixel) + '.png')
        plt.savefig(image_path)
        images.append(image_path)
        pp.savefig()

# Plot intensity dependent JV curves
    i = 0
    i_max = len(group_by_label_pixel_HL) - 1
    j = 0
    for ng_HL, ng_LH in zip(group_by_label_pixel_HL, group_by_label_pixel_LH):
        name_HL = ng_HL[0]
        group_HL = ng_HL[1]
        name_LH = ng_LH[0]
        group_LH = ng_LH[1]
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]
        jsc_max = max(max(group_HL['Jsc']), max(group_LH['Jsc']))
        c_div = 1 / len(group_HL)
        if i % 4 == 0:
            plt.figure(figsize=(A4_width, A4_height), dpi=300)
            plt.suptitle('Intensity dependent JV curves',
                         fontsize=10,
                         fontweight='bold')
            j += 1
        plt.subplot(2, 2, (i % 4) + 1)
        k = 0
        for path_HL, path_LH, intensity_HL, intensity_LH in zip(
                group_HL['File_Path'], group_LH['File_Path'],
                group_HL['Intensity'], group_LH['Intensity']):
            data_HL = np.genfromtxt(path_HL, delimiter='\t')
            data_LH = np.genfromtxt(path_LH, delimiter='\t')
            data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
            data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
            plt.plot(data_HL[:, 0],
                     data_HL[:, 1],
                     c=cmap(k * c_div),
                     label=str(round(intensity_HL * 100, 1)) + ' W/cm^2')
            plt.plot(data_LH[:, 0], data_LH[:, 1], c=cmap(k * c_div))
            k += 1
        plt.legend(loc='best', fontsize=7)
        plt.xlabel('Applied voltage (V)', fontsize=9)
        plt.ylabel('J (mA/cm^2)', fontsize=9)
        plt.ylim([-jsc_max * 1.05, jsc_max * 1.05])
        plt.title(
            str(label) + ', pixel ' + str(pixel) + ', ' + str(variable) + ', '
            + str(value),
            fontsize=8)
        if (i % 4 == 3) or (i == i_max):
            plt.tight_layout()
            plt.subplots_adjust(top=0.92)
            image_path = log_file_jv.replace(
                '.txt', '_intensity_JV_' + str(j) + '.png')
            plt.savefig(image_path)
            images.append(image_path)
            pp.savefig()
        i += 1

# Plot eqe graphs if experiment data exists
if os.path.exists(folderpath + folderpath_eqe + filepath_eqe):
    data = pd.read_csv(folderpath + folderpath_eqe + filepath_eqe,
                       delimiter='\t',
                       header=0,
                       names=['Label', 'Pixel', 'Variable', 'Value',
                              'Position', 'Int_Jsc', 'Mismatch', 'Area',
                              'Frequency', 'File_Path'])
    sorted_data_eqe = data.sort_values(
        ['Variable', 'Value', 'Label', 'Pixel', 'Int_Jsc'],
        ascending=[True, True, True, True, False])
    sorted_data_eqe = sorted_data_eqe.drop_duplicates(['Label', 'Pixel'])
    filtered_data_HL_temp = pd.merge(filtered_data_HL,
                                     sorted_data_eqe,
                                     on=['Label', 'Pixel'],
                                     how='inner')
    filtered_data_LH_temp = pd.merge(filtered_data_LH,
                                     sorted_data_eqe,
                                     on=['Label', 'Pixel'],
                                     how='inner')
    sorted_data_eqe_HL = sorted_data_eqe
    sorted_data_eqe_LH = sorted_data_eqe
    sorted_data_eqe_HL['SS_Jsc'] = list(filtered_data_HL_temp['Jsc'])
    sorted_data_eqe_LH['SS_Jsc'] = list(filtered_data_LH_temp['Jsc'])
    grouped_by_label = sorted_data_eqe.groupby(['Label'])

    i = 0
    i_max = len(sorted_data_eqe.drop_duplicates(['Label'])) - 1
    j = 0
    for name, group in grouped_by_label:
        label = group['Label'].unique()[0]
        variable = group['Variable'].unique()[0]
        value = group['Value'].unique()[0]
        c_div = 1 / len(group)

        if i % 4 == 0:
            plt.figure(figsize=(A4_width, A4_height), dpi=300)
            plt.suptitle('External quantum efficiency',
                         fontsize=10,
                         fontweight='bold')
            j += 1
        plt.subplot(2, 2, (i % 4) + 1)
        k = 0
        for path, pixel in zip(group['File_Path'], group['Pixel']):
            data = np.genfromtxt(path, delimiter='\t')
            plt.plot(data[:, 0],
                     data[:, 1],
                     c=cmap(k * c_div),
                     label=str(pixel))
            k += 1
        plt.legend(loc='lower center', fontsize=7)
        plt.xlabel('Wavelength (nm)', fontsize=9)
        plt.ylabel('EQE (%)', fontsize=9)
        plt.xlim([np.min(data[:, 0]), np.max(data[:, 0])])
        plt.ylim([0, 100])
        plt.title(
            str(label) + ', ' + str(variable) + ', ' + str(value),
            fontsize=8)
        if (i % 4 == 3) or (i == i_max):
            plt.tight_layout()
            plt.subplots_adjust(top=0.92)
            image_path = log_file_jv.replace('.txt', '_EQE_' + str(j) + '.png')
            plt.savefig(image_path)
            images.append(image_path)
            pp.savefig()
        i += 1

# Plot graph of measured Jsc against integrated Jsc
    marker = itertools.cycle((',', 'D', 'x', 'o', '*', '^'))
    color = itertools.cycle(
        ('black', 'blue', 'red', 'green', 'purple', 'magenta', 'cyan'))

    plt.figure(figsize=(A4_width, A4_height), dpi=300)
    plt.suptitle('Measured vs. Integrated Jsc', fontsize=10, fontweight='bold')
    zipped = zip(sorted_data_eqe_HL['Label'], sorted_data_eqe_HL['Pixel'],
                 sorted_data_eqe_HL['SS_Jsc'], sorted_data_eqe_HL['Int_Jsc'],
                 sorted_data_eqe_LH['Label'], sorted_data_eqe_LH['Pixel'],
                 sorted_data_eqe_LH['SS_Jsc'], sorted_data_eqe_LH['Int_Jsc'])
    for label_HL, pixel_HL, ss_jsc_HL, int_jsc_HL, label_LH, pixel_LH, ss_jsc_LH, int_jsc_LH in zipped:
        plt.scatter(
            ss_jsc_HL,
            int_jsc_HL,
            label=(str(label_HL) + ', pixel ' + str(pixel_HL) + ', H->L'),
            marker=next(marker),
            s=30,
            c=next(color))
        plt.scatter(
            ss_jsc_LH,
            int_jsc_LH,
            label=(str(label_LH) + ', pixel ' + str(pixel_LH) + ', L->H'),
            marker=next(marker),
            s=30,
            c=next(color))
    plt.plot(sorted_data_eqe_HL['SS_Jsc'],
             sorted_data_eqe_HL['SS_Jsc'],
             c='black',
             label='Int Jsc = SS Jsc')
    plt.legend(loc='upper left', scatterpoints=1, fontsize=9)
    plt.xlabel('Solar simulator Jsc (mA/cm^2)', fontsize=9)
    plt.ylabel('Integrated Jsc from EQE (mA/cm^2)', fontsize=9)
    image_path = log_file_jv.replace('.txt', '_EQE_Jscs' + '.png')
    plt.savefig(image_path)
    images.append(image_path)
    pp.savefig()

# Get weather data and add to dataframe
weather_data_path = r'C:/SolarSimData/WeatherData/WxLog.csv'
weather_data = pd.read_csv(
    weather_data_path,
    delimiter=',',
    skiprows=5,
    header=0,
    usecols=[0, 1, 11, 13, 14, 16, 17, 19, 20, 22],
    parse_dates=[[0, 1]],
    names=['Date', 'time', 'Lab T', 'Lab RH', 'Fume hood T', 'Fume hood RH',
           'Solar sim T', 'Solar sim RH', 'Dessicator T', 'Dessicator RH'])

# Filter dataframe to leave only data from last 7 days
one_week_ago = date.today() - timedelta(days=7)
weather_data = weather_data.loc[weather_data.Date_time > one_week_ago]

# Plot weather data and save images
fig = plt.figure(figsize=(A4_width, A4_height), dpi=300)
fig.suptitle('Weather report for last 7 days', fontsize=10, fontweight='bold')
gs = gridspec.GridSpec(2, 1)
ax1 = fig.add_subplot(gs[0, 0])
ax1.plot_date(weather_data['Date_time'],
              weather_data['Lab T'],
              fmt='-',
              label='Lab')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Fume hood T'],
              fmt='-',
              label='Fumehood')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Solar sim T'],
              fmt='-',
              label='Solar sim')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Dessicator T'],
              fmt='-',
              label='Dessicator')
ax1.set_xticklabels([])
ax1.set_ylim([15, 30])
ax1.legend(loc='best')
ax1.set_ylabel('Temperature (C)', fontsize=8)
ax2 = fig.add_subplot(gs[1, 0])
ax2.plot_date(weather_data['Date_time'],
              weather_data['Lab RH'],
              fmt='-',
              label='Lab')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Fume hood RH'],
              fmt='-',
              label='Fumehood')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Solar sim RH'],
              fmt='-',
              label='Solar sim')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Dessicator RH'],
              fmt='-',
              label='Dessicator')
ax2.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%y %H:%M'))
ax2.set_ylim([0, 100])
ax2.legend(loc='best')
ax2.set_ylabel('Relative humidity (%)', fontsize=8)
plt.xticks(rotation='35', ha='right', fontsize=8)
gs.update(wspace=0.3, hspace=0.05)
image_path = log_file_jv.replace('.txt', '_weather_' + '.png')
plt.savefig(image_path)
images.append(image_path)
pp.savefig()

# Close pdf of saved figures
pp.close()

# Create a powerpoint presentation to add figures to
prs = Presentation()

# Add title page with experiment title, date, and username
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = experiment_title
subtitle.text = exp_date + ', ' + username

# Add slide with table for manual completion of experimental details
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes
rows = 17
cols = 6
left = Inches(0.15)
top = Inches(0.02)
width = prs.slide_width - Inches(0.25)
height = prs.slide_height - Inches(0.05)
table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(0.8)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.0)
table.columns[4].width = Inches(2.0)
table.columns[5].width = Inches(1.7)

# write column headings
table.cell(0, 0).text = 'Label'
table.cell(0, 1).text = 'Substrate'
table.cell(0, 2).text = 'Bottom contact'
table.cell(0, 3).text = 'Perovskite'
table.cell(0, 4).text = 'Top contact'
table.cell(0, 5).text = 'Top electrode'

# fill in label column
i = 1
for item in sorted(sorted_data['Label'].unique()):
    table.cell(i, 0).text = str(item)
    i += 1

# Add images to blank slides
blank_slide_layout = prs.slide_layouts[6]
height = prs.slide_height
width = prs.slide_width
left = top = Inches(0)
for image in images:
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(image, left, top, height=height)

# Save powerpoint presentation
prs.save(log_file_jv.replace('.txt', '_summary.pptx'))
