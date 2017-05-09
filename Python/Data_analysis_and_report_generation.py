# This script takes the measurement log files and loads them as Pandas
# DataFrames for manipulation and then plotting.

import argparse
import itertools
import os
from datetime import date, timedelta

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import scipy as sp
from matplotlib import axes
from pptx import Presentation
from pptx.util import Inches
from scipy import constants

import reportgenlib as rgl

# Bind subplot formatting methods in reportgenlib to the matplotlib.axes.Axes
# class.
axes.Axes.set_axes_props = rgl.set_axes_props
axes.Axes.subboxplot = rgl.subboxplot
axes.Axes.subbarchart = rgl.subbarchart

# Parse folder path and log file name from command line arguments.
# Remember to include the "python" keyword before the call to the python file
# from the command line, e.g. python example.py "arg1" "arg2". Folder paths
# must use forward slashes to separate subfolders.
parser = argparse.ArgumentParser(description='Process data files')
parser.add_argument(
    'folder_path',
    metavar='folder_path',
    type=str,
    help='Absolute path to the folder containing all of the experimental data')
parser.add_argument('log_file_name',
                    metavar='log_file_name',
                    type=str,
                    help='Name of the master J-V log file')
args = parser.parse_args()

# Define folder and file paths
folderpath = args.folder_path
folderpath_jv = 'J-V/'
folderpath_time = 'Time Dependence/'
folderpath_maxp = 'Max P Stabilisation/'
folderpath_intensity = 'Intensity Dependence/'
folderpath_eqe = 'EQE/'
filepath_jv = args.log_file_name
filepath_eqe = r'_EQE_LOG.txt'
log_file_jv = folderpath + folderpath_jv + filepath_jv
log_file_time = folderpath + folderpath_time + filepath_jv
log_file_maxp = folderpath + folderpath_maxp + filepath_jv
log_file_intensity = folderpath + folderpath_intensity + filepath_jv
log_file_eqe = folderpath + folderpath_eqe + filepath_jv

# Create folders for storing files generated during analysis
analysis_folder = folderpath + 'Analysis/'
image_folder = analysis_folder + 'Figures/'
if os.path.exists(analysis_folder):
    pass
else:
    os.makedirs(analysis_folder)
if os.path.exists(image_folder):
    pass
else:
    os.makedirs(image_folder)

# Get username, date, and experiment title from folderpath for the title page
# of the powerpoint version of the report.
folderpath_split = folderpath.split('/')
username = folderpath_split[2]
date_title_split = folderpath_split[5].split(' ')
exp_date = date_title_split[0]
experiment_title = ' '.join(date_title_split[1:])

# Set physical constants
kB = constants.Boltzmann
q = constants.elementary_charge
c = constants.speed_of_light
h = constants.Planck
T = 300

# Read in data from JV log file
data = pd.read_csv(log_file_jv,
                   delimiter='\t',
                   header=0,
                   names=['Label', 'Pixel', 'Condition', 'Variable', 'Value',
                          'Position', 'Jsc', 'Voc', 'PCE', 'FF', 'Area',
                          'Stabil_Level', 'Stabil_time', 'Meas_delay', 'Vmp',
                          'File_Path', 'Scan_rate', 'Scan_direction',
                          'Intensity'])

# Read scan numbers from file paths and add scan number column
# to dataframe
scan_num = []
for path in data['File_Path']:
    scan_i = path.find('scan', len(path) - 12)
    scan_num.append(path[scan_i:].strip('scan').strip('.txt'))
data['scan_num'] = pd.Series(scan_num, index=data.index)

# Create a powerpoint presentation to add figures to.
prs = Presentation()

# Add title page with experiment title, date, and username.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = experiment_title
subtitle.text = exp_date + ', ' + username

# Add slide with table for manual completion of experimental details.
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes
rows = 17
cols = 6
left = Inches(0.15)
top = Inches(0.02)
width = prs.slide_width - Inches(0.25)
height = prs.slide_height - Inches(0.05)
table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(0.8)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.0)
table.columns[4].width = Inches(2.0)
table.columns[5].width = Inches(1.7)

# write column headings
table.cell(0, 0).text = 'Label'
table.cell(0, 1).text = 'Substrate'
table.cell(0, 2).text = 'Bottom contact'
table.cell(0, 3).text = 'Perovskite'
table.cell(0, 4).text = 'Top contact'
table.cell(0, 5).text = 'Top electrode'

# Define dimensions used for adding images to slides
A4_height = 7.5
A4_width = 10
height = prs.slide_height * 0.95 / 2
width = prs.slide_width * 0.95 / 2

# Create dictionaries that define where to put images on slides
# in the powerpoint presentation.
lefts = {'0': Inches(0),
         '1': prs.slide_width - width,
         '2': Inches(0),
         '3': prs.slide_width - width}
tops = {'0': prs.slide_height * 0.05,
        '1': prs.slide_height * 0.05,
        '2': prs.slide_height - height,
        '3': prs.slide_height - height}

# Use the extra analysis function to perform extra analysis and create new
# series for the dataframe
Rs_grad = []
Rsh_grad = []
for i in range(len(data['File_Path'])):
    file = data['File_Path'][i]
    new_path = file.replace('\\', '\\\\')
    Jsc = data['Jsc'][i]
    Vmp = data['Vmp'][i]
    Voc = data['Voc'][i]
    Area = data['Area'][i]
    if data['Condition'][i] == 'Light':
        extra_parameters = rgl.extra_JV_analysis(new_path, Jsc, Vmp, Voc, Area)
        Rs_grad.append(extra_parameters[0])
        Rsh_grad.append(extra_parameters[1])
    else:
        Rs_grad.append(0)
        Rsh_grad.append(0)

# Add new series to the dataframe
data['Rs_grad'] = pd.Series(Rs_grad, index=data.index)
data['Rsh_grad'] = pd.Series(Rsh_grad, index=data.index)

# Sort data
sorted_data = data.sort_values(['Variable', 'Value', 'Label', 'Pixel', 'PCE'],
                               ascending=[True, True, True, True, False])

# Fill in label column of device info table in ppt
i = 1
for item in sorted(sorted_data['Label'].unique()):
    table.cell(i, 0).text = str(item)
    i += 1

# Filter data
filtered_data = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01)]
filtered_data_HL = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01) &
                               (sorted_data.Scan_direction == 'HL')]
filtered_data_LH = sorted_data[(sorted_data.Condition == 'Light') & (
    sorted_data.FF > 0.1) & (sorted_data.FF < 0.9) & (sorted_data.Jsc > 0.01) &
                               (sorted_data.Scan_direction == 'LH')]
filtered_data = filtered_data.drop_duplicates(['Label', 'Pixel'])
filtered_data_HL = filtered_data_HL.drop_duplicates(['Label', 'Pixel'])
filtered_data_LH = filtered_data_LH.drop_duplicates(['Label', 'Pixel'])

# Drop pixels only working in one scan direction.
# First get the inner merge of the Label and Pixel columns, i.e. drop rows
# where label and pixel combination only occurs in one scan direction.
filtered_data_HL_t = filtered_data_HL[['Label', 'Pixel']].merge(
    filtered_data_LH[['Label', 'Pixel']],
    on=['Label', 'Pixel'],
    how='inner')
filtered_data_LH_t = filtered_data_LH[['Label', 'Pixel']].merge(
    filtered_data_HL[['Label', 'Pixel']],
    on=['Label', 'Pixel'],
    how='inner')

# Then perform inner merge of full filtered data frames with the merged
# label and pixel dataframes to get back all pixel data that works in both
# scan directions
filtered_data_HL = filtered_data_HL.merge(filtered_data_HL_t,
                                          on=['Label', 'Pixel'],
                                          how='inner')
filtered_data_LH = filtered_data_LH.merge(filtered_data_LH_t,
                                          on=['Label', 'Pixel'],
                                          how='inner')

# Calculate proportion of working pixels.
# Create groups of all and working pixels from complete and filtered
# dataframes.
group_var_s = sorted_data.drop_duplicates(['Label', 'Pixel'])
group_var_s = group_var_s.groupby(['Variable'])
group_var_f = filtered_data_HL.drop_duplicates(['Label', 'Pixel'])
group_var_f = group_var_f.groupby(['Variable'])

# For each variable, if there are any working pixels for that variable
# calculate the yield for each value of the variable if there are any working
# pixels for that value. Else add zeros. Also, make a list of names
# for each variable to use as the x-axis on bar charts.
yields_var = []
names_yield_var = []
for var_key in list(group_var_s.groups.keys()):
    group_val_s = group_var_s.get_group(var_key).groupby(['Value'])
    if var_key in list(group_var_f.groups.keys()):
        group_val_f = group_var_f.get_group(var_key).groupby(['Value'])
        yields_val = []
        names_yield_val = []
        for val_key in list(group_val_s.groups.keys()):
            names_yield_val.append(val_key)
            if val_key in list(group_val_f.groups.keys()):
                yields_val.append(len(group_val_f.get_group(val_key)) * 100 /
                                  len(group_val_s.get_group(val_key)))
            else:
                yields_val.append(0)
        yields_var.append(yields_val)
        names_yield_var.append(names_yield_val)
    else:
        yields_var.append([0] * len(group_val_s))
        names_yield_var.append(list(group_val_s.groups.keys()))

# For each variable, if there are any working pixels for that variable
# calculate the yield for each label if there are any working
# pixels for that label. Else add zeros. Also, make a list of names
# for each variable to use as the x-axis on bar charts.
yields_var_lab = []
names_yield_var_lab = []
for var_key in list(group_var_s.groups.keys()):
    group_lab_s = group_var_s.get_group(var_key).groupby(['Label'])
    if var_key in list(group_var_f.groups.keys()):
        group_lab_f = group_var_f.get_group(var_key).groupby(['Label'])
        yields_lab = []
        names_yield_lab = []
        for lab_key in list(group_lab_s.groups.keys()):
            names_yield_lab.append(lab_key)
            if lab_key in list(group_lab_f.groups.keys()):
                yields_lab.append(len(group_lab_f.get_group(lab_key)) * 100 /
                                  8)
            else:
                yields_lab.append(0)
        yields_var_lab.append(yields_lab)
        names_yield_var_lab.append(names_yield_lab)
    else:
        yields_var_lab.append([0] * len(group_lab_s))
        names_yield_var_lab.append(list(group_lab_s.groups.keys()))

# To generate box plots the data needs to be grouped first by variable
grouped_by_var_HL = filtered_data_HL.groupby('Variable')
grouped_by_var_LH = filtered_data_LH.groupby('Variable')

# Then it needs to be grouped by variable value. Each of these groupings is
# appended to a list that is iterated upon later to generate the plots.
# Create variables holding a dictionary for accessing lists of data for
# boxplots.
boxplotdata_HL = rgl.boxplotdata(grouped_by_var_HL)
boxplotdata_LH = rgl.boxplotdata(grouped_by_var_LH)

# Iterate through the lists of grouped data to produce boxplots. Each plot
# will contain all data from all values of a variable including a 'Control'
# sample if one is given. Multiple plots are created if there is more than
# one variable.
for i in range(len(boxplotdata_HL['var_names'])):
    if boxplotdata_HL['var_names'][i] != 'Control':
        for j in range(len(boxplotdata_HL['names_var'])):
            if boxplotdata_HL['names_var'][j][0] == 'Control':
                boxplotdata_HL['names_var'][i].append(boxplotdata_HL[
                    'names_var'][j][0])
                boxplotdata_HL['Jsc_var'][i].append(boxplotdata_HL['Jsc_var'][
                    j][0])
                boxplotdata_HL['Voc_var'][i].append(boxplotdata_HL['Voc_var'][
                    j][0])
                boxplotdata_HL['FF_var'][i].append(boxplotdata_HL['FF_var'][j][
                    0])
                boxplotdata_HL['PCE_var'][i].append(boxplotdata_HL['PCE_var'][
                    j][0])
                boxplotdata_HL['Rs_var'][i].append(boxplotdata_HL['Rs_var'][j][
                    0])
                boxplotdata_HL['Rsh_var'][i].append(boxplotdata_HL['Rsh_var'][
                    j][0])

# Build a list of x values for scatter plots that will overlay
# boxplots.
        x = []
        for k in range(1, 1 + len(boxplotdata_HL['names_var'][i])):
            x.extend([k] * len(boxplotdata_HL['Jsc_var'][i][k - 1]))

# Create new slide and box plots for PV parameters, save figures,
# and add them to the new slide
        data_slide = rgl.title_image_slide(
            prs, boxplotdata_HL['var_names'][i] + ' basic parameters')
        params = ['Jsc', 'Voc', 'FF', 'PCE']
        for ix, p in enumerate(params):
            image_path = image_folder + 'boxplot_' + p + '.png'
            rgl.create_save_boxplot(p, boxplotdata_HL, boxplotdata_LH, i, x,
                                    image_path)
            data_slide.shapes.add_picture(image_path,
                                          left=lefts[str(ix)],
                                          top=tops[str(ix)],
                                          height=height)

# Create new slide, box plots, and bar charts for parameters, save
# figures and add them to the new slide
        data_slide = rgl.title_image_slide(
            prs, boxplotdata_HL['var_names'][i] +
            ' series and shunt resistances; yields')
        params = ['Rs', 'Rsh', 'yields_var', 'yields_var_lab']
        for ix, p in enumerate(params):
            if p.find('yield') == -1:
                image_path = image_folder + 'boxplot_' + p + '.png'
                rgl.create_save_boxplot(p, boxplotdata_HL, boxplotdata_LH, i,
                                        x, image_path)
            elif p == 'yields_var':
                image_path = image_folder + 'barchart_' + p + '.png'
                rgl.create_save_barchart(yields_var, names_yield_var, i,
                                         image_path)
            elif p == 'yields_var_lab':
                image_path = image_folder + 'barchart_' + p + '.png'
                rgl.create_save_barchart(yields_var_lab, names_yield_var_lab,
                                         i, image_path)
            data_slide.shapes.add_picture(image_path,
                                          left=lefts[str(ix)],
                                          top=tops[str(ix)],
                                          height=height)

# Clear all figures from memory
plt.close('all')

# Group data by label and sort ready to plot graph of all pixels per substrate
re_sort_data = filtered_data.sort_values(['Label', 'Pixel'],
                                         ascending=[True, True])
grouped_by_label = re_sort_data.groupby('Label')

# Define a colormap for JV plots
cmap = plt.cm.get_cmap('rainbow')

# Create lists of varibales, values, and labels for labelling figures
substrates = re_sort_data.drop_duplicates(['Label'])
variables = list(substrates['Variable'])
values = list(substrates['Value'])
labels = list(substrates['Label'])

# Create figures, save images and add them to powerpoint slide
i = 0
for name, group in grouped_by_label:

    # Create a new slide after every four graphs are produced
    if i % 4 == 0:
        data_slide = rgl.title_image_slide(
            prs, 'JV scans of every working pixel, page ' + str(int(i / 4)))

# Create figure, axes, y=0 line, and title
    fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
    ax = fig.add_subplot(1, 1, 1)
    ax.axhline(0, lw=0.5, c='black')
    ax.set_title(str(labels[i]) + ', ' + str(variables[i]) + ', ' + str(values[
        i]))

    c_div = 1 / len(group)
    pixels = list(group['Pixel'])
    max_group_jsc = max(list(group['Jsc']))

    # Import data for each pixel and plot on axes
    j = 0
    for file in group['File_Path']:

        if '_LH_' in new_path:
            data_LH_path = file
            data_HL_path = file.replace('_LH_', '_HL_')
        else:
            data_HL_path = file
            data_LH_path = file.replace('_HL_', '_LH_')
        data_LH = np.genfromtxt(data_LH_path, delimiter='\t')
        data_HL = np.genfromtxt(data_HL_path, delimiter='\t')
        data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
        data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]

        ax.plot(data_LH[:, 0],
                data_LH[:, 1],
                label=pixels[j],
                c=cmap(j * c_div),
                lw=2.0)
        ax.plot(data_HL[:, 0], data_HL[:, 1], c=cmap(j * c_div), lw=2.0)

        j += 1

# Format the axes
    ax.set_xlabel('Applied bias (V)')
    ax.set_ylabel('J (mA/cm^2)')
    ax.set_xlim([np.min(data_HL[:, 0]), np.max(data_HL[:, 0])])
    ax.set_ylim([-max_group_jsc * 1.1, max_group_jsc * 1.1])

    # Adjust plot width to add legend outside plot area
    box = ax.get_position()
    ax.set_position([box.x0, box.y0, box.width * 0.85, box.height])
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles, labels, loc='upper left', bbox_to_anchor=(1, 1))

    # Format the figure layout, save to file, and add to ppt
    image_path = image_folder + 'jv_all_' + str(labels[i]) + '.png'
    fig.savefig(image_path, bbox_extra_artists=(lgd, ), bbox_inches='tight')
    data_slide.shapes.add_picture(image_path,
                                  left=lefts[str(i % 4)],
                                  top=tops[str(i % 4)],
                                  height=height)

    # Close figure
    plt.close(fig)

    i += 1

# filter dataframe to leave on the best pixel for each variable value
sort_best_pixels = filtered_data.sort_values(['Variable', 'Value', 'PCE'],
                                             ascending=[True, True, False])
best_pixels = sort_best_pixels.drop_duplicates(['Variable', 'Value'])

# get parameters for defining position of figures in subplot, attempting to
# make it as square as possible
no_of_subplots = len(best_pixels['File_Path'])
subplot_rows = np.ceil(no_of_subplots**(1 / 2))
subplot_cols = np.ceil(no_of_subplots / subplot_rows)

# create lists of varibales and values for labelling figures
variables = list(best_pixels['Variable'])
values = list(best_pixels['Value'])
jscs = list(best_pixels['Jsc'])

# Loop for iterating through best pixels dataframe and picking out JV data
# files. Each plot contains forward and reverse sweeps, both light and dark.
i = 0
for file in best_pixels['File_Path']:

    # Create a new slide after every four graphs are produced
    if i % 4 == 0:
        data_slide = rgl.title_image_slide(
            prs, 'Best pixel JVs, page ' + str(int(i / 4)))

# Create figure, axes, y=0 line, and title
    fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
    ax = fig.add_subplot(1, 1, 1)
    ax.axhline(0, lw=0.5, c='black')
    ax.set_title(str(labels[i]) + ', ' + str(variables[i]) + ', ' + str(values[
        i]))

    # Import data for each pixel and plot on axes, ignoring errors. If
    # data in a file can't be plotted, just ignore it.
    if '_LH_' in new_path:
        JV_light_LH_path = file
        JV_light_HL_path = file.replace('_LH_', '_HL_')
    else:
        JV_light_HL_path = file
        JV_light_LH_path = file.replace('_HL_', '_LH_')

    try:
        JV_light_LH_data = np.genfromtxt(JV_light_LH_path, delimiter='\t')
        JV_light_HL_data = np.genfromtxt(JV_light_HL_path, delimiter='\t')
        JV_dark_LH_data = np.genfromtxt(
            JV_light_LH_path.replace('Light', 'Dark'),
            delimiter='\t')
        JV_dark_HL_data = np.genfromtxt(
            JV_light_HL_path.replace('Light', 'Dark'),
            delimiter='\t')
    except OSError:
        pass

    JV_light_LH_data = JV_light_LH_data[~np.isnan(JV_light_LH_data).any(
        axis=1)]
    JV_light_HL_data = JV_light_HL_data[~np.isnan(JV_light_HL_data).any(
        axis=1)]

    try:
        JV_dark_LH_data = JV_dark_LH_data[~np.isnan(JV_dark_LH_data).any(
            axis=1)]
        JV_dark_HL_data = JV_dark_HL_data[~np.isnan(JV_dark_HL_data).any(
            axis=1)]
    except NameError:
        pass

    ax.plot(JV_light_LH_data[:, 0],
            JV_light_LH_data[:, 1],
            label='L->H',
            c='red',
            lw=2.0)
    ax.plot(JV_light_HL_data[:, 0],
            JV_light_HL_data[:, 1],
            label='H->L',
            c='green',
            lw=2.0)
    try:
        ax.plot(JV_dark_LH_data[:, 0],
                JV_dark_LH_data[:, 1],
                label='L->H',
                c='blue',
                lw=2.0)
        ax.plot(JV_dark_HL_data[:, 0],
                JV_dark_HL_data[:, 1],
                label='H->L',
                c='orange',
                lw=2.0)
    except NameError:
        pass

# Format the axes
    ax.set_xlabel('Applied bias (V)')
    ax.set_ylabel('J (mA/cm^2)')
    ax.set_xlim(
        [np.min(JV_light_HL_data[:, 0]), np.max(JV_light_HL_data[:, 0])])
    ax.set_ylim([-jscs[i - 1] * 1.1, jscs[i - 1] * 1.1])
    ax.legend(loc='best')

    # Format the figure layout, save to file, and add to ppt
    image_path = image_folder + 'jv_best_' + str(variables[i]) + '_' + str(
        variables[i]) + '.png'
    fig.tight_layout()
    fig.savefig(image_path)
    data_slide.shapes.add_picture(image_path,
                                  left=lefts[str(i % 4)],
                                  top=tops[str(i % 4)],
                                  height=height)

    # Close figure
    plt.close(fig)

    i += 1

# Sort and filter data ready for plotting different scan rates/repeat scans
sorted_data_scan = data.sort_values(
    ['Variable', 'Value', 'Label', 'Pixel', 'scan_num'],
    ascending=[True, True, True, True, True])
filtered_scan_HL = sorted_data_scan[(sorted_data_scan.Condition == 'Light') & (
    sorted_data_scan.FF > 0.1) & (sorted_data_scan.FF < 0.9) & (
        sorted_data_scan.Jsc > 0.01) & (sorted_data_scan.Scan_direction == 'HL'
                                        )]
filtered_scan_LH = sorted_data_scan[(sorted_data_scan.Condition == 'Light') & (
    sorted_data_scan.FF > 0.1) & (sorted_data_scan.FF < 0.9) & (
        sorted_data_scan.Jsc > 0.01) & (sorted_data_scan.Scan_direction == 'LH'
                                        )]

# Drop pixels only working in one scan direction
filtered_data_HL = filtered_data_HL[filtered_data_HL.Label.isin(
    filtered_data_LH.Label.values) & filtered_data_HL.Pixel.isin(
        filtered_data_LH.Pixel.values)]
filtered_data_LH = filtered_data_LH[filtered_data_LH.Label.isin(
    filtered_data_HL.Label.values) & filtered_data_LH.Pixel.isin(
        filtered_data_HL.Pixel.values)]

# Create groups of data for each pixel for a given label
group_by_label_pixel_HL = filtered_scan_HL.groupby(['Label', 'Pixel'])
group_by_label_pixel_LH = filtered_scan_LH.groupby(['Label', 'Pixel'])

# Iterate through these groups and plot JV curves if more than one scan
# has been performed
i = 0
for iHL, iLH in zip(group_by_label_pixel_HL.indices,
                    group_by_label_pixel_LH.indices):
    group_HL = group_by_label_pixel_HL.get_group(iHL)
    group_LH = group_by_label_pixel_LH.get_group(iLH)

    if any(int(scan) > 0 for scan in group_HL['scan_num']):

        # Get label, variable, value, and pixel for title and image path
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]

        # Find maximum Jsc of the group for y-axis limits and number of
        # JV curves for division of the colormap for the curves
        jsc_max = max(max(group_HL['Jsc']), max(group_LH['Jsc']))
        c_div = 1 / len(group_HL)

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'Repeat scan/scan rate variation JV curves, page ' +
                str(int(i / 4)))

# Create figure, axes, y=0 line, and title
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.axhline(0, lw=0.5, c='black')
        ax.set_title(str(labels[i]) + ', ' + str(variables[i]) + ', ' + str(
            values[i]))

        # Open data files and plot a JV curve on the same axes for each scan
        j = 0
        for path_HL, path_LH, scan_rate_HL, scan_rate_LH, scan_num_HL, scan_num_LH in zip(
                group_HL['File_Path'], group_LH['File_Path'],
                group_HL['Scan_rate'], group_LH['Scan_rate'],
                group_HL['scan_num'], group_LH['scan_num']):

            data_HL = np.genfromtxt(path_HL, delimiter='\t')
            data_LH = np.genfromtxt(path_LH, delimiter='\t')
            data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
            data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]

            ax.plot(data_HL[:, 0],
                    data_HL[:, 1],
                    c=cmap(j * c_div),
                    label=str(scan_num_HL) + ', ' + str(scan_rate_HL) + ' V/s')
            ax.plot(data_LH[:, 0], data_LH[:, 1], c=cmap(j * c_div))

            j += 1

# Format the axes
        ax.set_xlabel('Applied bias (V)')
        ax.set_ylabel('J (mA/cm^2)')
        ax.set_xlim([np.min(data_HL[:, 0]), np.max(data_HL[:, 0])])
        ax.set_ylim([-jscs[i - 1] * 1.1, jscs[i - 1] * 1.1])

        # Adjust plot width to add legend outside plot area
        box = ax.get_position()
        ax.set_position([box.x0, box.y0, box.width * 0.8, box.height])
        handles, labels = ax.get_legend_handles_labels()
        lgd = ax.legend(handles,
                        labels,
                        loc='upper left',
                        bbox_to_anchor=(1, 1))

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'jv_repeats_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.savefig(image_path, bbox_extra_artists=(lgd), bbox_inches='tight')
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Close figure
        plt.close(fig)

        i += 1

# Build a time dependence log dataframe from file paths and J-V log file.
time_files = rgl.exp_file_list(folderpath + folderpath_time)
if time_files['exists']:
    time_files_df = rgl.build_log_df(time_files, sorted_data)

    sorted_time_files = time_files_df.sort_values(
        ['Label', 'Pixel', 'Scan_direction'],
        ascending=[True, True, True])
    sorted_time_files = sorted_time_files.drop_duplicates(['Label', 'Pixel'])

    i = 0
    for index, row in sorted_time_files.iterrows():

        # Get label, variable, value, and pixel for title and image path
        label = row['Label']
        variable = row['Variable']
        value = row['Value']
        pixel = row['Pixel']

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'J-t characterics, page ' + str(int(i / 4)))

# Open the data files
        path_HL = row['File_Path']
        path_LH = path_HL.replace('HL', 'LH')
        data_HL = np.genfromtxt(path_HL)
        data_LH = np.genfromtxt(path_LH)
        data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
        data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]
        data = np.vstack([data_HL, data_LH])

        # Create figure object
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)

        # Add axes for current density and format them
        ax1 = fig.add_subplot(2, 1, 1)
        ax1.plot(data[:, 0], data[:, 2], 'blue')
        yticks = rgl.calc_ticks(np.min(data[:, 2]), np.max(data[:, 2]))
        ax1.set_yticks(yticks)
        ax1.set_yticklabels(yticks)
        ax1.set_ylim([yticks[0], yticks[-1]])
        ax1.set_ylabel('J (mA/cm^2)')
        xticks = rgl.calc_ticks(np.min(data[:, 0]), np.max(data[:, 0]))
        ax1.set_xticks(xticks)
        ax1.set_xticklabels([])
        ax1.set_xlim([xticks[0], xticks[-1]])
        ax1.set_title(str(row['Label']) + ', pixel ' + str(row['Pixel']) + ', '
                      + str(row['Variable']) + ', ' + str(row['Value']))

        # Add axes for bias and format them
        ax2 = fig.add_subplot(2, 1, 2)
        ax2.plot(data[:, 0], data[:, 1], 'black')
        yticks = rgl.calc_ticks(np.min(data[:, 1]), np.max(data[:, 1]))
        ytick_range = yticks[-1] - yticks[0]
        ax2.set_yticks(yticks)
        ax2.set_yticklabels(yticks)
        ax2.set_ylim(
            [yticks[0] - 0.05 * ytick_range, yticks[-1] + 0.05 * ytick_range])
        ax2.set_ylabel('Bias (V)')
        ax2.set_xticks(xticks)
        ax2.set_xticklabels(xticks)
        ax2.set_xlim([xticks[0], xticks[-1]])
        ax2.set_xlabel('Time (s)')

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'jt_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Close figure
        plt.close(fig)

        i += 1

# Build a max power stabilisation log dataframe from file paths and J-V log
# file.
maxp_files = rgl.exp_file_list(folderpath + folderpath_maxp)
if maxp_files['exists']:
    maxp_files_df = rgl.build_log_df(maxp_files, sorted_data)

    i = 0
    for index, row in maxp_files_df.iterrows():

        # Get label, variable, value, and pixel for title and image path
        label = row['Label']
        variable = row['Variable']
        value = row['Value']
        pixel = row['Pixel']

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'Maximum power stabilisation, page ' + str(int(i / 4)))

# Open the data file
        path = row['File_Path']
        data = np.genfromtxt(path, delimiter='\t')
        data = data[~np.isnan(data).any(axis=1)]

        # Create figure object
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)

        # Add axes for bias and format them
        ax1 = fig.add_subplot(3, 1, 1)
        ax1.set_title(str(row['Label']) + ', pixel ' + str(row['Pixel']) + ', '
                      + str(row['Variable']) + ', ' + str(row['Value']))
        min_x = np.min(data[:, 0])
        max_x = np.max(data[:, 0])
        ax1.plot(data[:, 0], np.absolute(data[:, 1]), color='green')
        min_y = np.min(np.absolute(data[:, 1]))
        max_y = np.max(np.absolute(data[:, 1]))
        yrange = max_y - min_y
        ax1.set_ylim([min_y - yrange * 0.1, max_y + yrange * 0.1])
        ax1.set_ylabel('Bias (V)', color='green')
        ax1.set_xlim([min_x, max_x])
        ax1.set_xticklabels([])
        ax1.locator_params(axis='y', tight=False, nbins=4)

        # Add axes for J and format them
        ax2 = fig.add_subplot(3, 1, 2)
        ax2.plot(data[:, 0], np.absolute(data[:, 3]), color='red')
        min_y = np.min(np.absolute(data[:, 3]))
        max_y = np.max(np.absolute(data[:, 3]))
        yrange = max_y - min_y
        ax2.set_ylim([min_y - yrange * 0.1, max_y + yrange * 0.1])
        ax2.set_ylabel('|J| (mA/cm^2)', color='red')
        ax2.set_xlim([min_x, max_x])
        ax2.set_xticklabels([])
        ax2.locator_params(axis='y', tight=False, nbins=4)

        # Add axes for PCE and format them
        ax3 = fig.add_subplot(3, 1, 3)
        ax3.plot(data[:, 0], np.absolute(data[:, 5]), color='blue')
        min_y = np.min(np.absolute(data[:, 5]))
        max_y = np.max(np.absolute(data[:, 5]))
        yrange = max_y - min_y
        ax3.set_ylim([min_y - yrange * 0.1, max_y + yrange * 0.1])
        ax3.set_ylabel('PCE (%)', color='blue')
        ax3.set_xlim([min_x, max_x])
        ax3.set_xlabel('Time (s)')
        ax3.locator_params(axis='y', tight=False, nbins=4)

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'mppt_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.subplots_adjust(hspace=0.05)
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Close figure
        plt.close(fig)

        i += 1

# Plot inensity dependent graphs if experiment data exists
if os.path.exists(folderpath + folderpath_intensity + filepath_jv):
    data = pd.read_csv(folderpath + folderpath_intensity + filepath_jv,
                       delimiter='\t',
                       header=0,
                       names=['Label', 'Pixel', 'Condition', 'Variable',
                              'Value', 'Position', 'Jsc', 'Voc', 'PCE', 'FF',
                              'Area', 'Stabil_Level', 'Stabil_time',
                              'Meas_delay', 'Vmp', 'File_Path', 'Scan_rate',
                              'Scan_direction', 'Intensity'])

    # Sort, filter, and group intensity dependent data
    sorted_data_int = data.sort_values(['Label', 'Pixel', 'Intensity', 'PCE'],
                                       ascending=[True, True, True, False])
    filtered_data_int_HL = sorted_data_int[(
        sorted_data_int.Scan_direction == 'HL')]
    filtered_data_int_LH = sorted_data_int[(
        sorted_data_int.Scan_direction == 'LH')]
    filtered_data_int_HL = filtered_data_int_HL.drop_duplicates(
        ['Label', 'Pixel', 'Intensity'])
    filtered_data_int_LH = filtered_data_int_LH.drop_duplicates(
        ['Label', 'Pixel', 'Intensity'])
    group_by_label_pixel_HL = filtered_data_int_HL.groupby(['Label', 'Pixel'])
    group_by_label_pixel_LH = filtered_data_int_LH.groupby(['Label', 'Pixel'])

    # Plot intensity dependent JV curve parameters
    for ng_HL, ng_LH in zip(group_by_label_pixel_HL, group_by_label_pixel_LH):

        # Unpack group name and data
        name_HL = ng_HL[0]
        group_HL = ng_HL[1]
        name_LH = ng_LH[0]
        group_LH = ng_LH[1]

        # Get label, variable, value, and pixel for title and image path
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]

        # Perfom linear fit to intensity dependence of Jsc
        m_HL, c_HL, r_HL, p_HL, se_HL = sp.stats.linregress(
            group_HL['Intensity'] * 100, group_HL['Jsc'])
        m_LH, c_LH, r_LH, p_LH, se_LH = sp.stats.linregress(
            group_LH['Intensity'] * 100, group_LH['Jsc'])
        r_sq_HL = r_HL**2
        r_sq_LH = r_LH**2

        # Create new slide (do this every iteration of the loop because each
        # loop creates four graphs)
        data_slide = rgl.title_image_slide(
            prs, 'Intensity dependence ' + str(label) + ', ' + str(variable) +
            ', ' + str(value) + ', pixel ' + str(pixel))

        # Create intensity dependence of Jsc figure
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.scatter(group_HL['Intensity'] * 100,
                   group_HL['Jsc'],
                   c='blue',
                   label='H->L, ' + 'm=' + str(rgl.round_sig_fig(m_HL, 3)) +
                   ', c=' + str(rgl.round_sig_fig(c_HL, 3)) + ', R^2=' +
                   str(rgl.round_sig_fig(r_sq_HL, 3)))
        ax.scatter(group_LH['Intensity'] * 100,
                   group_LH['Jsc'],
                   c='red',
                   label='L->H, ' + 'm=' + str(rgl.round_sig_fig(m_LH, 3)) +
                   ', c=' + str(rgl.round_sig_fig(c_LH, 3)) + ', R^2=' +
                   str(rgl.round_sig_fig(r_sq_LH, 3)))

        # Adjust plot width to add legend outside plot area
        ax.legend(loc='upper left', scatterpoints=1, prop={'size': 9})

        # Plot linear fits
        ax.plot(group_HL['Intensity'] * 100,
                group_HL['Intensity'] * 100 * m_HL + c_HL,
                c='blue')
        ax.plot(group_LH['Intensity'] * 100,
                group_LH['Intensity'] * 100 * m_LH + c_LH,
                c='red')

        # Format axes
        ax.set_xlabel('Light intensity (mW/cm^2)')
        ax.set_ylabel('Jsc (mA/cm^2)')
        ax.set_xlim([0, np.max(group_HL['Intensity'] * 100) * 1.05])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'intensity_jsc_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(0)],
                                      top=tops[str(0)],
                                      height=height)

        # Close figure
        plt.close(fig)

        # Perfom linear fit to ln(Jsc) dependence of Voc to estimate n and
        # J0 assuming single diode equivalent circuit
        m_HL, c_HL, r_HL, p_HL, se_HL = sp.stats.linregress(
            np.log(group_HL['Jsc']), group_HL['Voc'])
        m_LH, c_LH, r_LH, p_LH, se_LH = sp.stats.linregress(
            np.log(group_LH['Jsc']), group_LH['Voc'])
        r_sq_HL = r_HL**2
        r_sq_LH = r_LH**2
        n_HL = m_HL * q / (kB * T)
        n_LH = m_LH * q / (kB * T)
        j0_HL = np.exp(-c_HL / m_HL)
        j0_LH = np.exp(-c_LH / m_LH)

        # Create ln(Jsc) dependence of Voc figure
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.scatter(
            np.log(group_HL['Jsc']),
            group_HL['Voc'],
            c='blue',
            label='H->L, ' + 'n=' + str(rgl.round_sig_fig(n_HL, 3)) + ', J_0='
            + ('%.2e' % j0_HL) + ' (mA/cm^2)' + ', R^2=' +
            str(rgl.round_sig_fig(r_sq_HL, 3)))
        ax.scatter(
            np.log(group_LH['Jsc']),
            group_LH['Voc'],
            c='red',
            label='L->H, ' + 'n=' + str(rgl.round_sig_fig(n_LH, 3)) + ', J_0='
            + ('%.2e' % j0_LH) + ' (mA/cm^2)' + ', R^2=' +
            str(rgl.round_sig_fig(r_sq_LH, 3)))

        # Adjust plot width to add legend outside plot area
        ax.legend(loc='upper left', scatterpoints=1, prop={'size': 9})

        # Plot linear fits
        ax.plot(
            np.log(group_HL['Jsc']),
            np.log(group_HL['Jsc']) * m_HL + c_HL,
            c='blue')
        ax.plot(
            np.log(group_LH['Jsc']),
            np.log(group_LH['Jsc']) * m_LH + c_LH,
            c='red')

        # Format axes
        ax.set_xlabel('ln(Jsc) (mA/cm^2)')
        ax.set_ylabel('Voc (V)')
        ax.set_ylim(
            [np.min(group_HL['Voc']) * 0.95, np.max(group_HL['Voc']) * 1.05])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'intensity_voc_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(1)],
                                      top=tops[str(1)],
                                      height=height)

        # Close figure
        plt.close(fig)

        # Create intensity dependence of FF figure
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.scatter(group_HL['Intensity'] * 100,
                   group_HL['FF'],
                   c='blue',
                   label='H->L')
        ax.scatter(group_LH['Intensity'] * 100,
                   group_LH['FF'],
                   c='red',
                   label='L->H')
        ax.legend(loc='best', scatterpoints=1)
        ax.set_xlabel('Light intensity (mW/cm^2)')
        ax.set_ylabel('FF')
        ax.set_xlim([0, np.max(group_HL['Intensity'] * 100) * 1.05])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'intensity_ff_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(2)],
                                      top=tops[str(2)],
                                      height=height)

        # Close figure
        plt.close(fig)

        # Create intensity dependence of PCE figure
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.scatter(group_HL['Intensity'] * 100,
                   group_HL['PCE'],
                   c='blue',
                   label='H->L')
        ax.scatter(group_LH['Intensity'] * 100,
                   group_LH['PCE'],
                   c='red',
                   label='L->H')
        ax.legend(loc='best', scatterpoints=1)
        ax.set_xlabel('Light intensity (mW/cm^2)')
        ax.set_ylabel('PCE (%)')
        ax.set_xlim([0, np.max(group_HL['Intensity'] * 100) * 1.05])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'intensity_pce_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(3)],
                                      top=tops[str(3)],
                                      height=height)

        # Close figure
        plt.close(fig)

# Plot intensity dependent JV curves
    i = 0
    for ng_HL, ng_LH in zip(group_by_label_pixel_HL, group_by_label_pixel_LH):

        # Unpack group name and data
        name_HL = ng_HL[0]
        group_HL = ng_HL[1]
        name_LH = ng_LH[0]
        group_LH = ng_LH[1]

        # Get label, variable, value, and pixel for title and image path
        label = group_HL['Label'].unique()[0]
        variable = group_HL['Variable'].unique()[0]
        value = group_HL['Value'].unique()[0]
        pixel = group_HL['Pixel'].unique()[0]

        # Find maximum Jsc of the group for y-axis limits and number of
        # JV curves for division of the colormap for the curves
        jsc_max = max(max(group_HL['Jsc']), max(group_LH['Jsc']))
        c_div = 1 / len(group_HL)

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'Intensity dependent JV curves, page ' + str(int(i / 4)))

# Create figure, axes, y=0 line, and title
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.axhline(0, lw=0.5, c='black')
        ax.set_title(str(label) + ', pixel ' + str(pixel) + ', ' + str(
            variable) + ', ' + str(value))

        # Open data files and plot a JV curve on the same axes for each scan
        j = 0
        for path_HL, path_LH, intensity_HL, intensity_LH in zip(
                group_HL['File_Path'], group_LH['File_Path'],
                group_HL['Intensity'], group_LH['Intensity']):

            data_HL = np.genfromtxt(path_HL, delimiter='\t')
            data_LH = np.genfromtxt(path_LH, delimiter='\t')
            data_HL = data_HL[~np.isnan(data_HL).any(axis=1)]
            data_LH = data_LH[~np.isnan(data_LH).any(axis=1)]

            ax.plot(data_HL[:, 0],
                    data_HL[:, 1],
                    c=cmap(j * c_div),
                    label=str(round(intensity_HL * 100, 1)) + ' mW/cm^2')
            ax.plot(data_LH[:, 0], data_LH[:, 1], c=cmap(j * c_div))

            j += 1

# Format axes
        ax.set_xlabel('Applied voltage (V)')
        ax.set_ylabel('J (mA/cm^2)')
        ax.set_xlim([np.min(data_LH[:, 0]), np.max(data_LH[:, 0])])
        ax.set_ylim([-jsc_max * 1.05, jsc_max * 1.05])

        # Adjust plot width to add legend outside plot area
        box = ax.get_position()
        ax.set_position([box.x0, box.y0, box.width * 0.7, box.height])
        handles, labels = ax.get_legend_handles_labels()
        lgd = ax.legend(handles,
                        labels,
                        loc='upper left',
                        bbox_to_anchor=(1, 1),
                        prop={'size': 9})

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'jv_intensity_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '_' + str(pixel) + '.png'
        fig.savefig(image_path,
                    bbox_extra_artists=(lgd, ),
                    bbox_inches='tight')
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Close figure
        plt.close(fig)

        i += 1

# Plot eqe graphs if experiment data exists
if os.path.exists(folderpath + folderpath_eqe + filepath_eqe):
    data = pd.read_csv(folderpath + folderpath_eqe + filepath_eqe,
                       delimiter='\t',
                       header=0,
                       names=['Label', 'Pixel', 'Variable', 'Value',
                              'Position', 'Int_Jsc', 'Mismatch', 'Area',
                              'Frequency', 'File_Path'])

    # Sort, filter, and group data
    sorted_data_eqe = data.sort_values(
        ['Variable', 'Value', 'Label', 'Pixel', 'Int_Jsc'],
        ascending=[True, True, True, True, False])
    sorted_data_eqe = sorted_data_eqe.drop_duplicates(['Label', 'Pixel'])
    filtered_data_HL_temp = pd.merge(filtered_data_HL,
                                     sorted_data_eqe,
                                     on=['Label', 'Pixel'],
                                     how='inner')
    filtered_data_LH_temp = pd.merge(filtered_data_LH,
                                     sorted_data_eqe,
                                     on=['Label', 'Pixel'],
                                     how='inner')
    sorted_data_eqe_HL = sorted_data_eqe
    sorted_data_eqe_LH = sorted_data_eqe
    sorted_data_eqe_HL['SS_Jsc'] = list(filtered_data_HL_temp['Jsc'])
    sorted_data_eqe_LH['SS_Jsc'] = list(filtered_data_LH_temp['Jsc'])
    grouped_by_label = sorted_data_eqe.groupby(['Label'])

    # Plot EQE graphs with all pixels on the same device on the same plot
    i = 0
    for name, group in grouped_by_label:

        # Get label, variable, and value for title and image path
        label = group['Label'].unique()[0]
        variable = group['Variable'].unique()[0]
        value = group['Value'].unique()[0]

        # Get colormap increment
        c_div = 1 / len(group)

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'External quantum efficiency, page ' + str(int(i / 4)))

# Create figure, axes, and title
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.set_title(str(label) + ', ' + str(variable) + ', ' + str(value))

        # Plot EQE spectra for each pixel on same plot
        j = 0
        for path, pixel in zip(group['File_Path'], group['Pixel']):
            data = np.genfromtxt(path, delimiter='\t')
            ax.plot(data[:, 0],
                    data[:, 1],
                    c=cmap(j * c_div),
                    label=str(pixel))
            j += 1

# Format axes
        ax.legend(loc='best')
        ax.set_xlabel('Wavelength (nm)')
        ax.set_ylabel('EQE (%)')
        ax.set_xlim([np.min(data[:, 0]), np.max(data[:, 0])])
        ax.set_ylim([0, 100])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'eqe_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Clear figures from memory
        plt.close(fig)

        i += 1

# Plot spectral responsivity graphs with all pixels on the same device on
# the same plot
    i = 0
    for name, group in grouped_by_label:

        # Get label, variable, and value for title and image path
        label = group['Label'].unique()[0]
        variable = group['Variable'].unique()[0]
        value = group['Value'].unique()[0]

        # Get colormap increment
        c_div = 1 / len(group)

        # Start a new slide after every 4th figure
        if i % 4 == 0:
            data_slide = rgl.title_image_slide(
                prs, 'Spectral responsivity, page ' + str(int(i / 4)))

# Create figure, axes, and title
        fig = plt.figure(figsize=(A4_width / 2, A4_height / 2), dpi=300)
        ax = fig.add_subplot(1, 1, 1)
        ax.set_title(str(label) + ', ' + str(variable) + ', ' + str(value))

        # Plot EQE spectra for each pixel on same plot
        j = 0
        for path, pixel in zip(group['File_Path'], group['Pixel']):
            data = np.genfromtxt(path, delimiter='\t')
            ax.plot(data[:, 0],
                    data[:, 1] * data[:, 0] * 1e-9 * q / (100 * h * c),
                    c=cmap(j * c_div),
                    label=str(pixel))
            j += 1

# Format axes
        ax.legend(loc='best')
        ax.set_xlabel('Wavelength (nm)')
        ax.set_ylabel('Responsivity (A/W)')
        ax.set_xlim([np.min(data[:, 0]), np.max(data[:, 0])])
        ax.set_ylim([0, np.max(data[:, 0]) * 1e-9 * q / (h * c)])

        # Format the figure layout, save to file, and add to ppt
        image_path = image_folder + 'responsivity_' + str(label) + '_' + str(
            variable) + '_' + str(value) + '.png'
        fig.tight_layout()
        fig.savefig(image_path)
        data_slide.shapes.add_picture(image_path,
                                      left=lefts[str(i % 4)],
                                      top=tops[str(i % 4)],
                                      height=height)

        # Clear figures from memory
        plt.close(fig)

        i += 1

# Create markers and colours for measured Jsc against integrated Jsc
    marker = itertools.cycle((',', 'D', 'x', 'o', '*', '^'))
    color = itertools.cycle(
        ('black', 'blue', 'red', 'green', 'purple', 'magenta', 'cyan'))

    # Create ppt for measured vs. integrated Jsc scatter plot
    data_slide = rgl.title_image_slide(prs, 'Measured vs. Integrated Jsc')

    # Create figure and axes
    fig = plt.figure(figsize=(A4_width, A4_height), dpi=300)
    ax = fig.add_subplot(1, 1, 1)

    # Zip data to be unpacked in loop
    zipped = zip(sorted_data_eqe_HL['Label'], sorted_data_eqe_HL['Pixel'],
                 sorted_data_eqe_HL['SS_Jsc'], sorted_data_eqe_HL['Int_Jsc'],
                 sorted_data_eqe_LH['Label'], sorted_data_eqe_LH['Pixel'],
                 sorted_data_eqe_LH['SS_Jsc'], sorted_data_eqe_LH['Int_Jsc'])

    # Add 1 point for every measured-integrated pair per iteration
    for label_HL, pixel_HL, ss_jsc_HL, int_jsc_HL, label_LH, pixel_LH, ss_jsc_LH, int_jsc_LH in zipped:
        ax.scatter(
            ss_jsc_HL,
            int_jsc_HL,
            label=(str(label_HL) + ', pixel ' + str(pixel_HL) + ', H->L'),
            marker=next(marker),
            s=30,
            c=next(color))
        ax.scatter(
            ss_jsc_LH,
            int_jsc_LH,
            label=(str(label_LH) + ', pixel ' + str(pixel_LH) + ', L->H'),
            marker=next(marker),
            s=30,
            c=next(color))

# Plot line to indicate where measured Jsc = integrated Jsc
    ax.plot(sorted_data_eqe_HL['SS_Jsc'],
            sorted_data_eqe_HL['SS_Jsc'],
            c='black',
            label='Int Jsc = SS Jsc')

    # Format axes
    ax.set_xlabel('Solar simulator Jsc (mA/cm^2)')
    ax.set_ylabel('Integrated Jsc from EQE (mA/cm^2)')

    # Adjust plot width to add legend outside plot area
    box = ax.get_position()
    ax.set_position([box.x0, box.y0, box.width * 0.8, box.height])
    handles, labels = ax.get_legend_handles_labels()
    lgd = ax.legend(handles,
                    labels,
                    loc='upper left',
                    bbox_to_anchor=(1.01, 1),
                    scatterpoints=1,
                    prop={'size': 9})

    # Format the figure layout, save to file, and add to ppt
    image_path = image_folder + 'eqe_integrated_vs_ss' + '.png'
    fig.savefig(image_path, bbox_extra_artists=(lgd, ), bbox_inches='tight')
    data_slide.shapes.add_picture(image_path,
                                  left=lefts[str(0)],
                                  top=tops[str(0)],
                                  height=height * 2)

    # Close figure
    plt.close(fig)

# Get weather data and add to dataframe
weather_data_path = r'C:/SolarSimData/WeatherData/WxLog.csv'
weather_data = pd.read_csv(
    weather_data_path,
    delimiter=',',
    skiprows=5,
    header=0,
    usecols=[0, 1, 11, 13, 14, 16, 17, 19, 20, 22],
    parse_dates=[[0, 1]],
    names=['Date', 'time', 'Lab T', 'Lab RH', 'Fume hood T', 'Fume hood RH',
           'Solar sim T', 'Solar sim RH', 'Dessicator T', 'Dessicator RH'])

# Filter dataframe to leave only data from last 7 days
one_week_ago = date.today() - timedelta(days=7)
weather_data = weather_data.loc[weather_data.Date_time > one_week_ago]

# Create new slide for weather report
data_slide = rgl.title_image_slide(prs, 'Weather report for the last 7 days')

# Create figure
fig = plt.figure(figsize=(A4_width, A4_height), dpi=300)

# Add axes and data for temperature subplot then format
ax1 = fig.add_subplot(2, 1, 1)
ax1.plot_date(weather_data['Date_time'],
              weather_data['Lab T'],
              fmt='-',
              label='Lab')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Fume hood T'],
              fmt='-',
              label='Fumehood')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Solar sim T'],
              fmt='-',
              label='Solar sim')
ax1.plot_date(weather_data['Date_time'],
              weather_data['Dessicator T'],
              fmt='-',
              label='Dessicator')
ax1.set_xticklabels([])
ax1.set_ylim([15, 30])
ax1.legend(loc='best')
ax1.set_ylabel('Temperature (C)')

# Add axes and data for humidity subplot then format
ax2 = fig.add_subplot(2, 1, 2)
ax2.plot_date(weather_data['Date_time'],
              weather_data['Lab RH'],
              fmt='-',
              label='Lab')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Fume hood RH'],
              fmt='-',
              label='Fumehood')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Solar sim RH'],
              fmt='-',
              label='Solar sim')
ax2.plot_date(weather_data['Date_time'],
              weather_data['Dessicator RH'],
              fmt='-',
              label='Dessicator')
ax2.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%y %H:%M'))
ax2.set_ylim([0, 100])
ax2.legend(loc='best')
ax2.set_ylabel('Relative humidity (%)')
plt.xticks(rotation='40', ha='right')

# Format the figure layout, save to file, and add to ppt
image_path = image_folder + 'weather_report.png'
fig.tight_layout()
fig.savefig(image_path)
data_slide.shapes.add_picture(image_path,
                              left=lefts[str(0)],
                              top=tops[str(0)],
                              height=height * 2)

# Close figure
plt.close(fig)

# Save powerpoint presentation
prs.save(analysis_folder + filepath_jv.strip('LOG.txt') + 'summary.pptx')
